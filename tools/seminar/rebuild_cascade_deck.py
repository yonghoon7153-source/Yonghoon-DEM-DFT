#!/usr/bin/env python3
"""rebuild_cascade_deck.py — 세미나 덱의 **내용 영역**을 도형 다이어그램에서 실제 그림으로 바꾼다.

무엇을 하나
  Codex 판 덱(`Research_Seminar_2026_08_cascade_human_story.pptx`)은 머리말(제목·용어 띠·
  구분선·불릿·푸터)까지는 좋은데, 내용 영역이 전부 **사각형 + 좌측 액센트 스트라이프 +
  번호 체브론**으로 채워져 있다. 그게 AI 가 만든 티가 나는 부분이고, 정작 우리 계산 그림은
  25장 중 6장에만 들어 있다.

  이 도구는 슬라이드마다 머리말(shape 0–13)은 **손대지 않고**, 내용 영역만 비운 뒤
  지정한 그림을 앉힌다. 그림이 없는 슬라이드는 텍스트만 남겨 깨끗하게 둔다.

이 도구가 **못 하는 것**
  · 슬라이드를 새로 만들거나 순서를 바꾸지 않는다 (`add_slide.py` / `p:sldIdLst` 담당).
  · 머리말 레이아웃이 다른 템플릿에는 못 쓴다 — shape 0–13 이 머리말이라는 것이 전제다.
    다른 덱에 쓰려면 HEADER_N 을 다시 재야 한다.
  · 그림이 슬라이드 주제에 맞는지 판단하지 않는다. SPEC 은 사람이 채운다.
  · 렌더링 확인을 대신하지 않는다 — 이 환경 LibreOffice 가 이 덱을 못 여는 것과 별개로,
    최종 확인은 PowerPoint 에서 사람이 해야 한다.
"""
import copy, re, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

#: 머리말 도형 개수 — 이 인덱스 미만은 절대 건드리지 않는다.
HEADER_N = 14
#: ⚠ 표지(1장)는 머리말 배치가 다르다 — 푸터가 21–23 이라 clear_zone 을 쓰면 푸터가 날아간다.
#:    그래서 표지는 인덱스가 아니라 **y 띠**로 체브론만 지운다.
COVER = 1
COVER_BAND = (3.85, 5.15)
#: 내용 영역 (인치)
ZONE = dict(x=0.55, y=2.50, w=8.90, h=3.28)   # 아래 0.3 in 은 `< 그림이름 >` 자리
CAP_Y = 6.46
MUT = RGBColor(0x6b, 0x7d, 0x8f)
FIG = "/home/user/Yonghoon-DEM-DFT/"


#: 불릿·목록 색 규칙 — 한계·못한 것 = 빨강, 강조·우리 기여 = 파랑
RED = RGBColor(0xbe, 0x12, 0x3c)
BLUE = RGBColor(0x25, 0x63, 0xeb)
_MARK = re.compile(r"\[(r|b)\](.*?)\[/\1\]", re.S)


def _split_marked(text):
    """`[r]...[/r]` / `[b]...[/b]` 를 (조각, 색) 목록으로. 마크가 없으면 통째로 하나."""
    out, i = [], 0
    for m in _MARK.finditer(text):
        if m.start() > i:
            out.append((text[i:m.start()], None))
        out.append((m.group(2), RED if m.group(1) == "r" else BLUE))
        i = m.end()
    if i < len(text):
        out.append((text[i:], None))
    return out or [(text, None)]


def set_text(shape, s):
    """첫 run 의 서식을 유지한 채 텍스트 교체. 색 마크업이 있으면 run 을 쪼갠다.

    `text_frame.text=` 는 서식을 통째로 날리므로 쓰지 않는다.
    """
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return False
    proto = p0.runs[0]
    parts = _split_marked(s)
    proto.text = parts[0][0]
    if parts[0][1] is not None:
        proto.font.color.rgb = parts[0][1]
        proto.font.bold = True
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for txt, col in parts[1:]:
        nr = copy.deepcopy(proto._r)
        p0._p.append(nr)
        run = p0.runs[-1]
        run.text = txt
        if col is not None:
            run.font.color.rgb = col
            run.font.bold = True
        else:
            run.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            run.font.bold = False
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    return True


#: Arial 메트릭 호환 폰트 — 폭 측정용 (LibreOffice 가 이 환경에서 pptx 를 못 열어 렌더 QA 가 안 된다)
_FONT_PATHS = ("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
               "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")


def _text_pt(text, pt):
    """text 를 pt 크기로 그렸을 때의 폭(pt). Arial↔Liberation 은 메트릭 호환."""
    from PIL import ImageFont
    for fp in _FONT_PATHS:
        try:
            f = ImageFont.truetype(fp, 200)
            return f.getlength(text) / 200.0 * pt
        except Exception:
            continue
    return len(text) * pt * 0.5          # 폰트가 없으면 보수적 근사


def autofit_line(shape, floor_pt=15.0, step=0.75):
    """상자 한 줄에 들어갈 때까지 글자 크기를 줄인다. 줄여도 안 되면 floor 에서 멈춘다.

    ⚠ 문구는 건드리지 않는다 — 넘치면 **작아지지 삭제되지 않는다.**
    """
    tf = shape.text_frame
    runs = [r for pgh in tf.paragraphs for r in pgh.runs]
    if not runs:
        return None
    text = "".join(r.text for r in runs)          # 마크업은 이미 run 으로 풀려 있다
    box_pt = shape.width / 914400 * 72 - 6          # 좌우 안쪽 여백
    cur = (runs[0].font.size.pt if runs[0].font.size else 18.0)
    size = cur
    while size > floor_pt and _text_pt(text, size) > box_pt:
        size -= step
    if size < cur:
        from pptx.util import Pt as _Pt
        for r in runs:
            r.font.size = _Pt(size)
        return (cur, size, _text_pt(text, size) <= box_pt)
    return None


def clear_zone(slide):
    """머리말 뒤 도형을 전부 지운다."""
    shapes = list(slide.shapes)
    for sh in shapes[HEADER_N:]:
        sh._element.getparent().remove(sh._element)


def fit(pw, ph, bx, by, bw, bh):
    """(px,py,pw,ph) 를 박스 안에 비율 유지로 중앙 정렬."""
    s = min(bw / pw, bh / ph)
    w, h = pw * s, ph * s
    return bx + (bw - w) / 2, by + (bh - h) / 2, w, h


def add_pics(slide, paths, zone=ZONE, gap=0.22):
    """1–2장을 내용 영역에 앉힌다. 2장이면 좌우 분할."""
    n = len(paths)
    cell_w = (zone["w"] - gap * (n - 1)) / n
    for i, rel in enumerate(paths):
        p = FIG + rel
        if not os.path.exists(p):
            print(f"    ⚠ 그림 없음: {rel}")
            continue
        with Image.open(p) as im:
            iw, ih = im.size
        bx = zone["x"] + i * (cell_w + gap)
        x, y, w, h = fit(iw, ih, bx, zone["y"], cell_w, zone["h"])
        slide.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))


def add_caption(slide, text, y=CAP_Y):
    tb = slide.shapes.add_textbox(Inches(ZONE["x"]), Inches(y), Inches(ZONE["w"]), Inches(0.44))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    f = r.font
    f.size, f.italic, f.color.rgb, f.name = Pt(9.5), True, MUT, "Arial"


def add_statement(slide, text, y=3.6, size=17.5, x=1.05, w=7.90):
    """그림이 없는 장 — 도형 대신 한 문장을 크게."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size, r.font.color.rgb, r.font.name = Pt(size), RGBColor(0x1f, 0x29, 0x37), "Arial"
        p.space_after = Pt(6)


#: 후보 명단 — 계열별. 색은 계열 구분용이지 판정이 아니다.
ROSTER = [
 ("OXIDES", 37, "#be123c",
  "Ag₂O  Al₂O₃  B₂O₃  BaO  CaO  CoO  Cr₂O₃  CrO₃  Cu₂O  Fe₂O₃  Ga₂O₃  Gd₂O₃  GeO₂\n"
  "HfO₂  In₂O₃  La₂O₃  Li₂O  MgO  MnO  MoO₃  Na₂O  Nb₂O₅  Nd₂O₃  NiO  Sb₂O₅  Sc₂O₃\n"
  "SiO₂  Sm₂O₃  SnO₂  SrO  Ta₂O₅  TiO₂  V₂O₅  WO₃  Y₂O₃  ZnO  ZrO₂"),
 ("CHLORIDES", 19, "#65a30d",
  "AlCl₃  BaCl₂  CaCl₂  CrCl₃  FeCl₃  GaCl₃  HfCl₄  LaCl₃  LiCl  MgCl₂\n"
  "NbCl₅  NdCl₃  ScCl₃  SmCl₃  SrCl₂  TaCl₅  TiCl₄  YCl₃  ZrCl₄"),
 ("SULFIDES", 11, "#c05621",
  "Al₂S₃  As₂S₃  CaS  Ga₂S₃  GeS₂  Li₂S  MgS  Na₂S  Sb₂S₃  SiS₂  SnS₂"),
 ("FLUORIDES", 10, "#0284c7",
  "AlF₃  CaF₂  LaF₃  LiF  MgF₂  NdF₃  ScF₃  TiF₄  YF₃  ZrF₄"),
 ("BROMIDES", 5, "#7c3aed", "AlBr₃  CaBr₂  LiBr  MgBr₂  ZrBr₄"),
 ("NITRIDES", 5, "#0d9488", "AlN  Ca₃N₂  GaN  Li₃N  Mg₃N₂"),
 ("IODIDES", 4, "#6b7280", "AlI₃  LiI  MgI₂  NaI"),
]


def add_roster_boxes(slide, x0=0.60, y0=2.50, w=8.80, h=3.95):
    """계열별 박스 — 옅은 채움 + 계열색 머리글. 부록 전용(전체 91종)."""
    from pptx.enum.shapes import MSO_SHAPE
    left = ROSTER[:4]
    right = ROSTER[4:]
    colw = (w - 0.30) / 2
    for ci, group in enumerate((left, right)):
        y = y0
        for name, n, colr, names in group:
            nl = names.count("\n") + 1
            bh = 0.34 + 0.205 * nl + 0.14
            bx = x0 + ci * (colw + 0.30)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(bx), Inches(y), Inches(colw), Inches(bh))
            box.adjustments[0] = 0.055
            rgb = RGBColor.from_string(colr.lstrip("#").upper())
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*[int(255 - 0.10 * (255 - c)) for c in rgb])
            box.line.color.rgb = rgb
            box.line.width = Pt(1.25)
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.14)
            tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.05)
            pgh = tf.paragraphs[0]
            r1 = pgh.add_run(); r1.text = f"{name}   {n}"
            r1.font.size, r1.font.bold, r1.font.name = Pt(12), True, "Arial"
            r1.font.color.rgb = rgb
            for line in names.split("\n"):
                q = tf.add_paragraph()
                rr = q.add_run(); rr.text = line
                rr.font.size, rr.font.name = Pt(10.5), "Arial"
                rr.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            y += bh + 0.16


def add_roster_compact(slide, y=4.92):
    """표 아래에 전체 명단을 작은 글씨 한 덩어리로. 표가 구조를, 이 줄이 실체를 준다."""
    body = "  ·  ".join(f"{n}" for _, _, _, n in [])
    tb = slide.shapes.add_textbox(Inches(0.62), Inches(y), Inches(8.76), Inches(1.65))
    tf = tb.text_frame; tf.word_wrap = True
    for k, (name, n, colr, names) in enumerate(ROSTER):
        pgh = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        r1 = pgh.add_run(); r1.text = f"{name}  "
        r1.font.size, r1.font.bold, r1.font.name = Pt(9.5), True, "Arial"
        r1.font.color.rgb = RGBColor.from_string(colr.lstrip("#").upper())
        r2 = pgh.add_run(); r2.text = names.replace("\n", "  ")
        r2.font.size, r2.font.name = Pt(9.5), "Arial"
        r2.font.color.rgb = RGBColor(0x4b, 0x55, 0x63)
        pgh.space_after = Pt(1.5)


def add_roster(slide):
    """계열 머리글 + 명단. 두 단으로 나눠 앉힌다 (표·도형 없이 글자만)."""
    left = ROSTER[:3]
    right = ROSTER[3:]
    for col, group in ((0.62, left), (5.35, right)):
        y = 2.52
        for name, n, colr, names in group:
            tb = slide.shapes.add_textbox(Inches(col), Inches(y), Inches(4.10), Inches(0.26))
            r = tb.text_frame.paragraphs[0].add_run(); r.text = f"{name}   {n}"
            r.font.size, r.font.bold, r.font.name = Pt(12.5), True, "Arial"
            r.font.color.rgb = RGBColor.from_string(colr.lstrip("#").upper())
            y += 0.28
            nlines = names.count("\n") + 1
            tb2 = slide.shapes.add_textbox(Inches(col), Inches(y), Inches(4.10), Inches(0.24 * nlines))
            tf2 = tb2.text_frame; tf2.word_wrap = True
            for k, line in enumerate(names.split("\n")):
                pgh = tf2.paragraphs[0] if k == 0 else tf2.add_paragraph()
                rr = pgh.add_run(); rr.text = line
                rr.font.size, rr.font.name = Pt(11), "Arial"
                rr.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            y += 0.24 * nlines + 0.30


def add_fig_note(slide, text, x, y, w, size=11.5):
    """그림에서 뺀 해설 한 줄. 그림 안에 글을 넣지 않는 대신 여기에 둔다."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.30))
    tf = tb.text_frame; tf.word_wrap = True
    pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
    for txt, col in _split_marked(text):
        r = pgh.add_run(); r.text = txt
        r.font.size, r.font.name = Pt(size), "Arial"
        r.font.color.rgb = col or RGBColor(0x4b, 0x55, 0x63)
        r.font.bold = col is not None


def add_fig_label(slide, text, x, y, w):
    """`< Arrhenius plot >` 꼴의 그림 이름표. 레퍼런스 덱(2026-06-15)의 관례."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
    tf = tb.text_frame; tf.word_wrap = False
    pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
    r = pgh.add_run(); r.text = f"< {text} >"
    r.font.size, r.font.name, r.font.bold = Pt(13), "Arial", False
    r.font.color.rgb = RGBColor(0x1f, 0x29, 0x37)


def add_table(slide, rows, x, y, w, size=11.5, first_col_w=None, head_fill="EEF2F7"):
    """레퍼런스 덱 양식의 표 — 머리행 옅은 채움 + 굵게, 본문은 얇은 가로줄만.

    `[r]…[/r]` / `[b]…[/b]` 색 마크업을 셀 안에서도 쓸 수 있다.
    """
    nr, nc = len(rows), len(rows[0])
    rh = max(0.225, 0.0235 * size)
    h = rh + 1.06 * rh * 0 + rh * (nr - 1) + 0.035
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    if first_col_w:
        tbl.columns[0].width = Inches(first_col_w)
        rest = (w - first_col_w) / (nc - 1)
        for c in range(1, nc):
            tbl.columns[c].width = Inches(rest)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(rh)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            pgh = tf.paragraphs[0]
            pgh.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            for txt, col in _split_marked(str(val)):
                r = pgh.add_run(); r.text = txt
                r.font.size, r.font.name = Pt(size), "Arial"
                r.font.bold = (ri == 0) or (col is not None)
                r.font.color.rgb = col or (RGBColor(0x1f, 0x29, 0x37) if ri == 0
                                           else RGBColor(0x37, 0x41, 0x51))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(head_fill if ri == 0 else "FFFFFF")
    return y + h


def add_lines(slide, rows, x=1.05, y=2.75, w=7.9, size=13, lead=0.34, bold_prefix=True):
    """도형 없는 목록 — `제목 :: 본문` 형태를 굵게/보통으로 나눠 찍는다."""
    for i, row in enumerate(rows):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + i * lead), Inches(w), Inches(lead))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        if bold_prefix and "::" in row:
            head, body = row.split("::", 1)
            hp = _split_marked(head.strip())
            r1 = p.add_run(); r1.text = hp[0][0] + "  "
            r1.font.size, r1.font.bold, r1.font.name = Pt(size), True, "Arial"
            r1.font.color.rgb = hp[0][1] or RGBColor(0x1f, 0x29, 0x37)
            rest = body.strip()
        else:
            rest = row
        for txt, col in _split_marked(rest):
            r = p.add_run(); r.text = txt
            r.font.size, r.font.name = Pt(size), "Arial"
            r.font.color.rgb = col or RGBColor(0x37, 0x41, 0x51)
            r.font.bold = col is not None


# ── 슬라이드별 지시 ────────────────────────────────────────────────────────────
# 흐름 (사용자 지정 8절):
#   ① 왜 스크리닝인가(2–3)  ② 남들은 어떻게(4–5)  ③ 후보군(6–7)
#   ④ 9단계 설계(8–16)      ⑤ 결과(17–20)         ⑥ 다음 판(21)
#   ⑦ ML·co-doping(22)      ⑧ discussion(23)      부록(24–27)
# 색: [r]…[/r] = 한계·못한 것,  [b]…[/b] = 강조·우리 기여
L = "litdb/figures/"
D = "docs/figures/"
S = "docs/figures/seminar/"
SPEC = {
 1:  dict(),   # 표지 — 4단 체브론만 걷어낸다 (부제가 이미 같은 말을 한다)

 # ① 왜 이걸 하나
 2:  dict(gloss='Argyrodite: the Li₆PS₅Cl crystal family  ·  Cathode-side: where the electrolyte is oxidised\nDendrite: metal growing through the electrolyte',
          src='Kang et al., Chem. Commun. 2026',
          title="The problem",
          header="A sulfide electrolyte fails in several ways at once",
          bullets=["Oxidation at the cathode, contact loss, cracking and dendrites are [r]coupled[/r], not separate.",
                   "So a modifier that helps one axis can [r]quietly cost you another[/r]."],
          pics=[L+"kang2026_intertwined_electrochemo_mechanical_sulfide_assb_review/fig_16.png"],
          cap="Kang et al., Chem. Commun. (2026), Fig. 16 — a literature map of how the failure modes feed each other. "
              "The three green boxes are the levers people actually pull. Not our result."),
 3:  dict(gloss='Lever: the variable you are allowed to change  ·  Charge recipe: how Li is added or removed to keep the cell neutral',
          src='Framing for this talk',
          title="The lever, and the cost",
          header="Doping the electrolyte is the lever we can screen cheaply",
          bullets=["Coating and anode work live at [r]interfaces[/r] — hard to build, harder to predict in advance.",
                   "Doping happens [b]inside the lattice[/b], so a computer can survey it before anyone synthesises."],
          lines=["The catch :: one element changes oxidation, Li transport and stiffness [r]all at once[/r] — several axes, not one.",
                 "The arithmetic :: [b]91 candidates[/b] × several sites × several charge recipes = thousands of structures.",
                 "So the question is not :: which is best? — it is [b]what deserves the expensive calculation next?[/b]"],
          lines_y=3.05, lines_size=14.5, lines_lead=0.66,
          statement="Screen broadly and cheaply.\nSpend the accurate methods only where the answer would change.",
          stmt_y=5.35, stmt_size=17.5),

 # ② 남들은 어떻게
 4:  dict(gloss='Gate: a pass / fail condition  ·  Fidelity: how accurate and how expensive the method is',
          src='Xiao 2019  ·  Kahle 2020',
          title="How others do it",
          header="Everyone narrows the space before spending on accuracy",
          bullets=["Xiao cuts 104,082 coating candidates to three with [b]sequential physical gates[/b].",
                   "Kahle instead [b]raises the accuracy[/b] stage by stage, ending in first-principles MD."],
          keep_pics=True,
          cap="Left: Xiao, Miara, Wang & Ceder, Joule 3 (2019) 1252, Fig. 1.   "
              "Right: Kahle, Marcolongo & Marzari, Energy Environ. Sci. 13 (2020) 928, Fig. 1."),
 5:  dict(gloss='Threshold: the numeric value of a gate  ·  Queue: the order in which expensive calculations get run',
          src='Sendek 2017',
          title="What we took",
          header="We borrowed the ordering, not the thresholds",
          bullets=["Their cut-offs were tuned for coatings; [r]a sulfide host would fail all of them[/r].",
                   "What transfers: [b]a cheap stage reorders the queue — it never certifies a material[/b]."],
          pics=[L+"sendek2017_ml_screening_12k_conductors/fig_1.png"],
          cap="Sendek, Yang, Cubuk, Duerloo, Cui & Reed, Energy Environ. Sci. 10 (2017) 306, Fig. 1 — physical gates first, "
              "then a classifier trained on only 40 measured conductors. Their honesty about small training sets is the part we reuse."),

 # ③ 후보군 — 부록에서 본문으로
 6:  dict(gloss='Family: compounds grouped by their dominant anion  ·  Roster: the full input list, before anything is judged',
          src='Input roster — not a shortlist',
          title="What we put in",
          header="91 compounds across seven families — chosen to span the chemistry",
          bullets=["Oxides for [b]strong bonds[/b]; halides because the host has Cl; nitrides and sulfides as controls.",
                   "⚠ The list leans toward [r]well-known stable compounds[/r] — that turns out to matter later."],
          table=dict(rows=[
                    ['family', 'n', 'why it is in the list', 'examples'],
                    ['oxides', '37', 'strong M–O bonds, expected to resist oxidation', 'Ag₂O  Al₂O₃  B₂O₃  BaO …'],
                    ['chlorides', '19', 'host already contains Cl', 'AlCl₃  BaCl₂  CaCl₂  CrCl₃ …'],
                    ['sulfides', '10', 'same anion as the host — control', 'Al₂S₃  CaS  Ga₂S₃  GeS₂ …'],
                    ['fluorides', '10', 'strongest bonds, worst Li mobility risk', 'AlF₃  CaF₂  LaF₃  LiF …'],
                    ['bromides', '5', 'larger halide — lattice-size control', 'AlBr₃  CaBr₂  LiBr  MgBr₂ …'],
                    ['nitrides', '5', 'high charge anion — charge-compensation control', 'AlN  Ca₃N₂  GaN  Li₃N …'],
                    ['iodides', '4', 'largest halide — mobility control', 'AlI₃  LiI  MgI₂  NaI'],
                 ], y=2.60, w=8.76, size=11.5, first_col_w=1.35),
          cap="Grouped by the dominant anion. Order inside a family carries no meaning — this is the input list, not a shortlist."),
 7:  dict(fig_label='candidate elements on the periodic table',
          fig_note='the small number is how many compounds of that element were run',
          gloss='Cation: the metal we substitute in  ·  Coverage: how many compounds of that element we ran',
          src='Our roster, drawn on the table',
          title="Where they sit",
          header="36 cation elements, from alkali metals to lanthanides",
          bullets=["The number in each box is [b]how many different compounds of that element[/b] we actually ran.",
                   "Coverage is broad in the transition metals and thin in the [r]heavy main group[/r]."],
          pics=[S+"roster_periodic_table.png"],
          cap="Our own roster. Colour marks the chemical family only — it is not a score and not a ranking."),

 # ④ 9단계
 8:  dict(fig_note='large cations take Li sites; small group-14 ions replace P',
          gloss='Sublattice: the set of equivalent positions one element occupies  ·  Aliovalent: a different charge from the ion it replaces',
          src='Our calculation  ·  site preference',
          title="STEP 1 — Where it sits",
          header="A dopant name is not yet a structure",
          bullets=["Li, P, S and Cl sublattices are distinct — an oxide needs [b]a cation site and an anion site[/b].",
                   "Aliovalent substitution adds or removes Li, and [r]more than one recipe balances the charge[/r]."],
          pics=[S+"step1_site_choice.png"],
          cap="Our calculation: which sublattice each cation prefers, plotted against its ionic radius. "
              "Big cations end up on Li sites; small group-14 ions replace P. Bars are the spread across our three runs."),
 9:  dict(fig_note='[r]the generator chose the site — it was not a variable we set[/r]',
          gloss='Configuration: one specific arrangement of atoms in the cell  ·  Generator: the code that enumerates candidate structures',
          src='Our calculation  ·  where the anion landed',
          title="STEP 2 — Building it",
          header="Every allowed site becomes its own candidate structure",
          bullets=["One compound produces [b]several structures[/b], not one — different sites, different charge recipes.",
                   "⚠ Which site wins is [r]an output of the generator[/r], not something we controlled."],
          pics=[S+"step2_anion_site.png"],
          cap="Our calculation: where the dopant's anion actually landed across the campaign. "
              "Three sublattices absorbed it in comparable numbers, so the site was never fixed by design."),
 10: dict(fig_note='[r]100 of 3,615 structures moved past 25 % and were dropped[/r]',
          gloss='Relaxation: moving atoms until the forces vanish  ·  Machine-learned potential: a fast stand-in for quantum-mechanical forces',
          src='Method note',
          title="STEP 3 — Cheap relaxation",
          header="Let a fast model decide which structures are even plausible",
          bullets=["A machine-learned potential relaxes in [b]minutes instead of hours[/b], so we can afford all of them.",
                   "Structures that collapse or refuse to converge drop out here."],
          pics=[S+"step3_survival.png"],
          fig_label="how far each structure moved",
          pic_zone=dict(x=0.55, y=2.48, w=8.90, h=2.20),
          statement="This stage answers one question: does the structure hold together?",
          stmt_y=5.42, stmt_size=15.5,
          ),
 11: dict(fig_note='[b]on this axis the three runs agree — the bars are short[/b]',
          gloss='Representative: the one structure carried forward  ·  Convergence: the optimiser actually finished',
          src='Our calculation  ·  stability band',
          title="STEP 4 — Picking one",
          header="One structure per candidate goes forward — and that is a choice",
          bullets=["We keep the converged, low-energy structure whose cell volume did not blow up.",
                   "⚠ Candidate counts differ per compound, so [r]this is not a ranking of the elements[/r]."],
          pics=[S+"step4_stability_band.png"],
          cap="Our calculation: relative stability against the undoped host, with the spread across our three runs. "
              "Here the runs agree, so a representative structure is defensible on this axis."),
 12: dict(fig_note='[b]every structure found a lower-energy arrangement once it could move[/b]',
          gloss='Anneal: a brief run at high temperature  ·  Local minimum: the nearest stable arrangement, not necessarily the right one',
          src='Method note',
          title="STEP 5 — Shaking it",
          header="A short heat-and-relax tests whether the structure survives",
          bullets=["Plain relaxation only finds the [r]nearest[/r] minimum; real synthesis explores much further.",
                   "So we heat briefly, then relax again — everything measured next depends on [b]bond lengths[/b]."],
          pics=[S+"step5_anneal_gain.png"],
          fig_label="energy gained by shaking the structure",
          pic_zone=dict(x=0.55, y=2.48, w=8.90, h=2.20),
          statement="500 K for 50 picoseconds, then relax again — long enough to cross a small barrier.",
          stmt_y=5.42, stmt_size=15.5,
          ),
 13: dict(fig_note='[r]valleys are low-energy regions, not verified channels[/r]',
          gloss='Bond-valence map: a cheap estimate of how comfortable an ion is at a point  ·  Percolation: a low-energy path that spans the crystal',
          src='Our calculation  ·  Li landscape',
          title="STEP 6 — The Li path",
          header="Map the energy landscape a lithium ion would have to cross",
          bullets=["We score every point for how comfortable a Li ion would be, then look for [b]connected valleys[/b].",
                   "We also count how much of the path the dopant itself [r]sits in the way of[/r]."],
          pics=[S+"step6_li_landscape.png"],
          cap="Our calculation: the static landscape, undoped versus B₂O₃-doped. "
              "⚠ Valleys are low-energy regions, not verified channels — this flags risk, it does not measure conductivity."),
 14: dict(fig_note='[r]a single case — not a pool-wide DFT validation[/r]',
          gloss='Equation of state: energy as the cell volume changes  ·  Stack pressure: the force holding the cell together',
          src='Our calculation  ·  checked against DFT',
          title="STEP 7 — Squeezing it",
          header="In a solid-state cell the mechanics are part of the performance",
          bullets=["Too stiff and particles [r]never make contact[/r]; too soft and the layer creeps.",
                   "We change the cell volume and apply small strains to get [b]stiffness and ductility[/b]."],
          pics=[S+"step7_equation_of_state.png"],
          cap="Our calculation: one equation-of-state case checked against DFT. "
              "⚠ A single-case check — the campaign-wide elastic numbers are model values compared within one convention."),
 15: dict(fig_note='[r]five candidates lose the window entirely — all late transition metals[/r]',
          gloss='Oxidation onset: the voltage at which decomposition becomes favourable  ·  Window: the voltage range where nothing is driven to decompose',
          src='Our calculation  ·  90 candidates',
          title="STEP 8 — The window",
          header="Raising the voltage pulls Li out — we compute where it breaks",
          bullets=["At each voltage we ask whether the material would rather stay whole or [r]split into other phases[/r].",
                   "We get the onset [b]and the products[/b] — an insulating one behaves nothing like a conducting one."],
          pics=[S+"step8_oxidation_windows.png"],
          fig_label="stability window of every candidate",
          cap="Our calculation: the voltage range over which each candidate is not driven to decompose. "
              "Late transition metals lose the window entirely. 0 K bulk thermodynamics — this says decomposition is allowed, not how fast."),
 16: dict(gloss='Proxy: something cheap that correlates with what you want  ·  Adhesion: how strongly two materials stick across an interface',
          src='Sundar et al., Adv. Sci. 2025',
          title="STEP 9 — What we skipped",
          header="Two calculations were designed in and never run",
          bullets=["[r]Real conductivity[/r] — long dynamics per candidate, days each. Step 6 stands in as a proxy.",
                   "[r]The interface[/r] — how the doped electrolyte meets the cathode. Not attempted at all."],
          pics=[L+"sundar2025_oxide_coating_screening_lpscl/fig_2.png"],
          cap="Sundar et al., Adv. Sci. (2025), Fig. 2 — the same oxide scored at four different interfaces. "
              "An element that looks safe against the electrolyte can look bad against the anode. We computed none of these four maps."),

 # ⑤ 결과
 17: dict(bullets=["Structures, relaxation, anneal, Li maps, mechanics and oxidation windows [b]all exist[/b].",
                   "Real conductivity, explicit interfaces and pool-wide DFT confirmation [r]do not[/r]."],
          gloss='Computed: a file-backed result exists  ·  Not computed: no result exists, and no proxy substitutes for it',
          src='Evidence inventory',
          title="What exists",
          header="Broad in structures, thin in the expensive evidence",
          table=dict(rows=[
                    ['what', 'how many of 270', 'status'],
                    ['candidate structures', '270', '[b]computed[/b]'],
                    ['fast relaxation', '270', '[b]computed[/b]'],
                    ['short anneal', '270', '[b]computed[/b]'],
                    ['static Li-path map', '270', '[b]computed[/b]'],
                    ['stiffness / EOS', '270', '[b]computed[/b]'],
                    ['oxidation window', '270', '[b]computed[/b]'],
                    ['conductivity from dynamics', '0', '[r]not run[/r]'],
                    ['cathode interface / adhesion', '0', '[r]not run[/r]'],
                 ], y=2.62, w=6.05, size=12, first_col_w=3.05),
          statement="A missing calculation\nstays missing.\n\nA proxy can reorder the queue;\nit cannot take the place\nof the answer.",
          stmt_x=6.95, stmt_y=2.95, stmt_size=14.5, stmt_w=2.55),
 18: dict(fig_note='[r]one species over three runs spans as much E as thirteen species[/r]',
          gloss='Exact formula: the atom counts in the simulated cell  ·  Placement: which site the dopant actually took',
          src='Our calculation  ·  label spread',
          title="Result 1 — Identity",
          header="The same dopant name did not mean the same material",
          bullets=["[b]59 of 90[/b] species kept one exact formula across our three runs; [r]31 did not[/r].",
                   "Where the site moved the value moved too — averaging the three [r]mixes different materials[/r]."],
          pics=[S+"result1_same_name_different_value.png"],
          fig_label="spread of the three runs, species by species",
          pic_zone=dict(x=0.55, y=2.50, w=8.90, h=1.98),
          table=dict(rows=[
                    ['axis', 'unit', 'within one species', 'across species', 'ratio', 'read as'],
                    ['relative stability', 'eV/atom', '0.037', '0.174', '0.21', '[b]chemistry dominates[/b]'],
                    ["Young's modulus", 'GPa', '6.553', '5.919', '1.11', '[r]placement dominates[/r]'],
                    ['Pugh ratio G/B', '—', '0.119', '0.200', '0.59', '[b]chemistry dominates[/b]'],
                    ['Li-path proxy', '—', '0.037', '0.013', '2.83', '[r]placement dominates[/r]'],
                 ], y=5.10, w=8.76, size=11, first_col_w=1.85),
          cap="Our calculation: each grey bar is one species, each dot one run. The bar is as tall as the gap between "
              "a dozen different species — which means placement, not chemistry, is driving that axis."),
 19: dict(gloss='Valence band: the highest filled electronic states  ·  Pinned: set by one element regardless of what else changes',
          src='Banik 2022  ·  our onsets',
          title="Result 2 — Oxidation",
          header="Sulfur sets the limit, and shifts away from it are conditional",
          bullets=["The top of the valence band is [b]sulfur[/b], so sulfur is oxidised first no matter what we add.",
                   "Late transition metals are the clear loss: [r]the window collapses[/r] rather than shifting."],
          pics=[L+"banik2022_substitutions_oxidative_stability_argyrodite/fig_4.png",
                S+"result2_onset_by_chemistry.png"],
          fig_label=["valence band edge of the host", "onset by dopant chemistry"],
          pic_zone=dict(x=0.55, y=2.48, w=8.90, h=1.70),
          table=dict(rows=[
                    ['dopant chemistry', 'n', 'median onset (V)', 'range (V)', 'windows lost'],
                    ['main-group', '21', '2.140', '1.72 – 2.36', '0'],
                    ['alk.earth', '17', '2.140', '2.06 – 2.14', '0'],
                    ['alkali', '10', '2.140', '2.06 – 2.14', '0'],
                    ['TM', '33', '2.024', '1.72 – 2.36', '[r]5[/r]'],
                    ['lanthanide', '9', '1.920', '1.89 – 2.12', '0'],
                 ], y=4.72, w=8.76, size=10.5, first_col_w=2.10),
          cap="Left: Banik et al. (2022), Fig. 4 — the valence-band edge of Li₆PS₅Cl is sulfur. "
              "Right: our onsets grouped by cation chemistry. Most sit pinned at the host value; we cannot yet attribute the exceptions to an element."),
 20: dict(fig_label='stability against blocked Li traffic',
          fig_note='[r]the more a dopant stabilises the lattice, the more Li traffic it blocks (r = −0.63)[/r]',
          gloss='Blocking fraction: how much of the Li path the dopant occupies  ·  Trade-off: a gain on one axis paid for on another',
          src='Our data  ·  Xiao 2019',
          title="Result 3 — The trade",
          header="What stabilises the lattice also blocks the traffic",
          bullets=["Across our candidates the two axes [r]pull against each other[/r] — the trend is real, not noise.",
                   "The same tension shows up in the literature on a completely different material set."],
          pics=[S+"result3_stability_vs_traffic.png",
                L+"xiao2019_cathode_coating_screening/fig_7.png"],
          cap="Left: our candidates — more stabilising means more Li sites blocked. "
              "Right: Xiao et al. (2019), Fig. 7 — across 411 oxides, more lithium means a lower oxidation limit. Same shape of problem."),

 # ⑥ 다음 판
 21: dict(gloss='Seed: a different random starting arrangement  ·  Boundary case: a candidate sitting near a decision threshold',
          src='Our calculation  ·  reference host',
          title="Next round",
          header="Fix the comparison before making it bigger",
          bullets=["[b]Freeze the site and the formula[/b], then repeat — that separates chemistry from placement.",
                   "Build real low concentrations in [b]larger cells[/b]; promote only boundary cases upward."],
          pics=[D+"bv_structure_panels_comp1.png"],
          cap="Our calculation: framework, Li percolation path and the bond-valence channel of the undoped host — "
              "the reference the next campaign is built on."),

 # ⑦ ML
 22: dict(gloss='Applicability domain: the region the training data actually covers  ·  Co-doping: two modifiers in the same lattice',
          src='Sendek 2017',
          title="ML & co-doping",
          header="Use the data to schedule calculations, not to invent results",
          bullets=["90 candidates gives [b]4,000 pairs[/b] — never computable, but predictable enough to prioritise.",
                   "⚠ On today's data a model would also learn our [r]placement noise[/r] — Result 1 comes first."],
          keep_pics=True,
          cap="Sendek et al. (2017), Fig. 4 — read the distance from the training data on the x-axis, not only the probability on the y-axis. "
              "Confidence outside the training domain is the failure mode to guard against."),

 # ⑧ discussion
 23: dict(bullets=["What we have is a [b]map of where the comparison is solid[/b] and where it wobbles.",
                   "What we do not have is a shortlist — and I would rather earn one than announce one."],
          gloss='Robust: survives being rebuilt a different way  ·  Validation budget: how many expensive calculations we can afford',
          src='Questions for the room',
          title="Discussion",
          header="Where I would like the room's opinion",
          lines=["1 · What first? :: A real concentration series with repeats, or the first explicit interface?",
                 "2 · What is “best”? :: The best single configuration, or the candidate that survives [b]being rebuilt[/b]?",
                 "3 · How many repeats? :: How many sites and seeds before a chemistry is worth promoting?",
                 "4 · Is a bulk-only screen worth handing to experiment? :: Or must the [r]interface[/r] be computed first?"],
          lines_y=2.90, lines_size=14.5, lines_lead=0.82,
          statement="Screen broadly  →  control the comparison  →  validate selectively",
          stmt_y=6.05, stmt_size=17.5),

 # 부록
 24: dict(gloss='Garnet: the Li₇La₃Zr₂O₁₂ oxide electrolyte family  ·  Survey: every dopant tried, not only the ones that worked',
          src='Anderson et al., 2024',
          title="Appendix — the twin",
          header="The same survey, done in a laboratory instead of a computer",
          bullets=["59 dopants synthesised and measured in a garnet — [b]our closest experimental twin[/b].",
                   "32 of our 36 cation elements appear in their table too."],
          keep_pics=True,
          cap="Anderson, Zolfaghar, Jonderian, Khaliullin & McCalla, Adv. Energy Mater. 14 (2024) 2304025, Fig. 1. "
              "Yellow marks dopants tried there for the first time."),
 25: dict(bullets=["The talk names each step by [b]the question it answers[/b]; the folders are named by what ran.",
                   "Two designed stages never produced evidence — they are marked [r]not run[/r]."],
          lines=["1 · Where it sits :: enumerate the sublattices and the charge recipes",
                 "2 · Building it :: generate one structure per allowed placement",
                 "3 · Cheap relaxation :: relax every candidate with the fast potential",
                 "4 · Picking one :: keep a converged, low-energy representative",
                 "5 · Shaking it :: brief anneal at 500 K, then relax again",
                 "6 · The Li path :: bond-valence landscape and the blocking count",
                 "7 · Squeezing it :: equation of state and elastic response",
                 "8 · The window :: grand-potential oxidation onset and products",
                 "9 · What we skipped :: [r]conductivity dynamics · explicit interface — never run[/r]"],
          lines_y=2.78, lines_size=13, lines_lead=0.385,
          gloss='Talk step: the question as told here  ·  Workflow stage: the folder that actually ran',
          src='Implementation map',
          title="Appendix — steps",
          header="How the nine talk steps map onto what actually ran"),
 26: dict(gloss='Replicate: the same thing measured again under the same conditions  ·  Quantised: rounded to the nearest whole atom',
          src='Label audit',
          title="Appendix — labels",
          header="Why our three runs are neither concentrations nor clean repeats",
          bullets=["All three collapsed to the [r]same substitution[/r] in a small cell — no concentration axis.",
                   "Regrouping by exact formula gives [b]59 same-formula[/b] and [r]31 changed-formula[/r] species."],
          lines=["What they are :: three directory names — x002, x005, x010.",
                 "What they actually were :: [r]the same integer substitution[/r] in a four-formula-unit cell.",
                 "What to do instead :: [b]fix the formula and the site[/b], then vary only the starting seed.",
                 "And separately :: change the cell size to create [b]a real concentration axis[/b]."],
          lines_y=2.95, lines_size=13.5, lines_lead=0.60,
          statement="A directory name is not a concentration,\nand three of them are not three repeats.",
          stmt_y=5.60, stmt_size=17),
 27: dict(bullets=["Left of the arrow is [b]what the method answers[/b]; right of it is what people assume it answers.",
                   "Every line on the right is a calculation we [r]have not run[/r]."],
          gloss='Proxy: cheap and correlated  ·  Validation: the matched, higher-fidelity calculation the claim actually needs',
          src='Method boundaries',
          title="Appendix — methods",
          header="Every method is named by the question it can actually answer",
          lines=["Fast relaxation :: does the structure hold together?  →  [r]not[/r] a formation energy",
                 "Short anneal :: does it survive being disturbed?  →  [r]not[/r] a conductivity",
                 "Bond-valence map :: is there a plausible Li path?  →  [r]not[/r] a diffusion coefficient",
                 "Equation of state, strain :: is the mechanical response plausible?  →  [r]not[/r] a measured modulus",
                 "Grand potential :: what can decompose, and into what?  →  [r]not[/r] a lifetime, and [r]not[/r] an interface"],
          lines_y=2.80, lines_size=13.5, lines_lead=0.58,
          statement="A proxy may decide what to compute next.\nIt may not stand in for the result.",
          stmt_y=5.80, stmt_size=17),
 29: dict(gloss="Site screening: choosing which sublattice a dopant occupies  ·  Coating: a separate compound placed between electrolyte and cathode",
          src="Miara 2015 via Anderson 2024 · Lee 2024 · Xiao 2019",
          pics=[S+"prior_site_defect_energy.png", S+"step1_site_choice.png"],
          fig_label=["45 dopants × 3 sites in a garnet (literature)", "91 compounds in our sulfide host (this work)"],
          fig_note=["[b]one cation at a time; the site is chosen by defect energy[/b]",
                    "[b]a whole compound enters — cation and anion together[/b]"],
          pic_zone=dict(x=0.55, y=2.50, w=8.90, h=2.20),
          table=dict(rows=[
                    ['study', 'host', 'what is varied', 'n', 'mechanical axis'],
                    ['Miara 2015 (via Anderson)', 'garnet oxide', 'one cation, 3 sites', '45', '[r]none[/r]'],
                    ['Lee 2024', 'argyrodite', 'elements already in the formula', '84', '[r]none[/r]'],
                    ['Xiao 2019', '[r]no host[/r] — free coating', 'the compound itself', '104,082', '[r]none[/r]'],
                    ['this work', 'argyrodite sulfide', '[b]a foreign binary compound[/b]', '91', '[b]included[/b]'],
                 ], y=5.12, w=8.76, size=10, first_col_w=2.30),
          cap="Left: Anderson et al. (2024), Fig. 3d — defect energies from Miara, Richards, Wang & Ceder, "
              "Chem. Mater. 27 (2015) 4040. Right: our own site-preference calculation."),
 28: dict(gloss="Roster: the full input list, before anything is judged  ·  Family: grouped by the dominant anion",
          src="master_batch roster · 91 compounds",
          roster_boxes=True,
          cap="All 91 planned compounds, grouped by family. Order inside a family carries no meaning."),
}

# ── 어체 통일 (2026-08-16) ────────────────────────────────────────────────────
# 모티브: 2026-01-13 연구세미나 덱 (Hong Rim Shin). 그 덱의 규칙은 셋이다.
#   · 제목 = 그 장에서 하는 일의 이름 (문장 아님, 동명사·명사구)
#   · ■ 헤더 = `주제 (n): 구체적 대상` — 번호를 붙여 계열을 만든다
#   · • 불릿 = 마침표로 끝나는 선언문. `To <목적>, <대상> was <동사>.` 가 기본형
# 아래 표가 SPEC 의 title/header/bullets 를 덮어쓴다.
VOICE = {
 2:  ("Failure modes in sulfide solid electrolytes",
      "Motivation (1): Electrochemical and mechanical degradation are coupled",
      ["Cathode-side oxidation, contact loss, cracking and dendrite growth [r]proceed together[/r].",
       "A modifier that improves one axis can [r]degrade another[/r] without appearing there."]),
 3:  ("Selecting the modification lever",
      "Motivation (2): Lattice doping can be screened before synthesis",
      ["Coating and anode engineering act at [r]interfaces[/r], which are hard to build and to predict.",
       "To survey a wide chemical space at low cost, [b]lattice doping was chosen[/b] as the variable."]),
 4:  ("Screening strategies in the literature",
      "Prior work (1): Broad spaces are narrowed before accurate methods",
      ["Xiao et al. reduced 104,082 coating candidates to three by [b]sequential physical gates[/b].",
       "Kahle et al. kept the candidates and instead [b]raised the level of theory[/b] stage by stage."]),
 5:  ("Adopted design principle",
      "Prior work (2): The ordering was adopted; the thresholds were not",
      ["Published thresholds were tuned for coatings, and [r]a sulfide host fails all of them[/r].",
       "The transferable rule: [b]a cheap stage reorders the queue[/b], it never certifies a material."]),
 6:  ("Candidate chemistry space",
      "Candidate set (1): 91 compounds across seven anion families",
      ["Oxides were included for [b]strong M–O bonding[/b]; halides because the host contains Cl.",
       "Sulfides, nitrides and heavy halides were included as [b]anion-family controls[/b]."]),
 7:  ("Elemental coverage of the candidate set",
      "Candidate set (2): 36 cation elements, alkali metals to lanthanides",
      ["Coverage is dense among transition metals and [r]sparse in the heavy main group[/r].",
       "The list is [r]biased toward well-characterised stable compounds[/r] — see Result 1."]),
 8:  ("Substitution site and charge compensation",
      "Step 1: The substitution site is enumerated before a structure is built",
      ["The Li, P, S and Cl sublattices are distinct; an oxide needs [b]a cation and an anion site[/b].",
       "Aliovalent substitution moves Li, and [r]more than one recipe balances the charge[/r]."]),
 9:  ("Candidate structure generation",
      "Step 2: Each allowed placement becomes a separate structure",
      ["To cover the placement space, [b]several structures per compound[/b] were generated.",
       "The winning site is [r]an output of the generator[/r], not a variable that was controlled."]),
 10: ("Low-cost structure relaxation",
      "Step 3: A machine-learned potential screens structures before DFT",
      ["To relax thousands of structures within budget, [b]a fast interatomic potential was used[/b].",
       "Energies are [r]relative within one convention[/r] — not formation energies or hull distances."]),
 11: ("Representative structure selection",
      "Step 4: One converged structure per candidate is carried forward",
      ["Structures were kept if they converged and the cell volume [b]stayed within 25 %[/b].",
       "Candidate counts differ per compound, so [r]this selection is not a ranking of elements[/r]."]),
 12: ("Thermal perturbation of the selected structure",
      "Step 5: A short anneal tests whether the arrangement survives",
      ["To escape the nearest local minimum, [b]each structure was heated briefly and relaxed again[/b].",
       "The trajectory is [r]not an equilibrium structure and not a conductivity measurement[/r]."]),
 13: ("Static lithium transport pathway",
      "Step 6: The Li energy landscape is mapped on the annealed geometry",
      ["To flag transport risk without dynamics, [b]a bond-valence landscape was computed[/b].",
       "Low-energy valleys are [r]structural pathways, not diffusion coefficients[/r]."]),
 14: ("Mechanical response of the doped lattice",
      "Step 7: Stiffness and compressibility are obtained from finite strains",
      ["To assess particle contact under stack pressure, [b]elastic moduli were computed[/b].",
       "The DFT comparison shown here is [r]a single case, not a pool-wide validation[/r]."]),
 15: ("Electrochemical stability window",
      "Step 8: The oxidation onset from a grand-potential construction",
      ["To locate decomposition, [b]the Li chemical potential was scanned as a voltage axis[/b].",
       "The result is [r]0 K bulk thermodynamics[/r] — not a rate, and not a passivation prediction."]),
 16: ("Calculations that were not performed",
      "Step 9: Two designed stages produced no evidence",
      ["[r]Conductivity from dynamics[/r] — days per candidate; Step 6 stands in as a proxy.",
       "[r]The cathode interface[/r] — not attempted; the figure shows what such a screen looks like."]),
 17: ("Inventory of available evidence",
      "Result summary: The screen is broad in structure and thin in transport",
      ["Structures, relaxation, anneal, Li maps, mechanics and oxidation windows [b]all exist[/b].",
       "Conductivity, explicit interfaces and pool-wide DFT confirmation [r]do not exist[/r]."]),
 18: ("Effect of structural identity on the measured value",
      "Result 1: The dopant name did not define a controlled comparison",
      ["Across the three runs, [b]59 of 90 species[/b] kept one exact formula and [r]31 did not[/r].",
       "Where the site moved the value moved with it, so [r]averaging mixes different materials[/r]."]),
 19: ("Dependence of oxidation onset on dopant chemistry",
      "Result 2: Sulfur sets the limit; shifts away from it are conditional",
      ["The valence-band edge is [b]sulfur[/b], so sulfur is oxidised first regardless of the dopant.",
       "Late transition metals are the clear loss: [r]the window collapses rather than shifting[/r]."]),
 20: ("Trade-off between stability and lithium transport",
      "Result 3: Lattice stabilisation and Li mobility oppose each other",
      ["Across our candidates the two axes [r]correlate negatively[/r] — the trend is not noise.",
       "The same trade-off appears in the literature for [b]an unrelated material set[/b]."]),
 21: ("Design of the next campaign",
      "Future work (1): Control is strengthened before scale is increased",
      ["To separate chemistry from placement, [b]the formula and site will be frozen[/b] before repeats.",
       "Real low concentrations require [b]larger cells[/b]; only boundary cases are promoted upward."]),
 22: ("Use of the data for machine learning",
      "Future work (2): The data schedules calculations, not replaces them",
      ["90 candidates give [b]about 4,000 pairs[/b] — not computable, but rankable for prioritisation.",
       "Trained on today's data a model would also learn our [r]placement noise[/r]; Result 1 comes first."]),
 23: ("Discussion",
      "Open questions for the group",
      ["What exists is [b]a map of where the comparison is reliable[/b], not a shortlist.",
       "The next decision is [b]which expensive calculation reduces uncertainty the most[/b]."]),
 24: ("Experimental counterpart in a garnet electrolyte",
      "Appendix A1: The same survey performed by synthesis",
      ["59 dopants were synthesised and measured — [b]the closest experimental analogue[/b].",
       "32 of our 36 cation elements also appear in that study."]),
 25: ("Mapping of talk steps onto the executed workflow",
      "Appendix A2: Nine questions, and where each one was answered",
      ["Each talk step is named by [b]the question it answers[/b]; folders are named by what ran.",
       "Two designed stages produced no evidence and are marked [r]not run[/r]."]),
 26: ("Interpretation of the three campaign labels",
      "Appendix A3: The labels are neither concentrations nor repeats",
      ["All three collapsed to the [r]same substitution[/r] in a small cell — no concentration axis.",
       "Regrouping by exact formula gives [b]59 same-formula[/b] and [r]31 changed-formula[/r] species."]),
 27: ("Boundaries of each computational method",
      "Appendix A4: Every method named by the question it can answer",
      ["The left column is [b]what the method answers[/b]; the right is what it is assumed to answer.",
       "Every entry on the right is a calculation that [r]has not been run[/r]."]),
}
VOICE[29] = ("Where prior screening stops",
             "Prior work (3): Site-resolved screening exists — for oxides, one cation",
             ["Ceder's group screened [b]45 cations across three garnet sites[/b] by defect energy in 2015.",
              "A coating is a [r]free-standing compound[/r]; a dopant must be [b]placed and charge-balanced[/b]."])

VOICE[28] = ("Full candidate roster",
             "Appendix A5: All 91 compounds, grouped by anion family",
             ["Every compound that entered the campaign is listed — [b]this is the denominator[/b].",
              "As₂S₃ was planned but [r]stopped during structure generation[/r]."])

for _n, (_t, _h, _b) in VOICE.items():
    SPEC.setdefault(_n, {}).update(title=_t, header=_h, bullets=_b)

# ── 선행연구 (3) 을 6번 자리에 끼운다 → 이후 장 번호가 한 칸씩 밀린다 ─────────
# 임시 키 29 로 써 놓고 여기서 6 으로 옮긴다. sldIdLst 재정렬(reordered29)과 짝이다.
_INSERT_AT, _TEMP_KEY = 6, 29
_shift = {}
for _k, _v in SPEC.items():
    if _k == _TEMP_KEY:
        _shift[_INSERT_AT] = _v
    elif _k >= _INSERT_AT:
        _shift[_k + 1] = _v
    else:
        _shift[_k] = _v
SPEC = _shift

# 텍스트만 손보고 내용 영역은 그대로 두는 장
KEEP_ZONE = set()


def main():
    src, dst = sys.argv[1], sys.argv[2]
    prs = Presentation(src)
    for i, slide in enumerate(prs.slides, 1):
        if i in KEEP_ZONE:
            sp_ = SPEC.get(i, {})
            if sp_.get("gloss"):
                set_text(slide.shapes[2], sp_["gloss"])
            if sp_.get("src"):
                set_text(slide.shapes[10], sp_["src"])
            if sp_.get("title"):
                set_text(slide.shapes[1], sp_["title"])
            if sp_.get("header"):
                set_text(slide.shapes[5], sp_["header"])
            for k, b in enumerate(sp_.get("bullets", [])):
                set_text(slide.shapes[7 + k * 2], b)
            for j, floor in ((1, 15.0), (5, 12.75)):
                autofit_line(slide.shapes[j], floor_pt=floor)
            print(f"  P{i:2d} 내용 유지 · 제목/헤더 갱신")
            continue
        spec = SPEC.get(i)
        if spec is None:
            # 지시 없는 장 — 좌측 스트라이프/체브론만 정리
            for sh in list(slide.shapes)[HEADER_N:]:
                w, h = sh.width / 914400, sh.height / 914400
                if w <= 0.12 and h >= 0.5:      # 좌측 액센트 스트라이프
                    sh._element.getparent().remove(sh._element)
            for j, floor in ((1, 15.0), (5, 12.75)):
                autofit_line(slide.shapes[j], floor_pt=floor)
            print(f"  P{i:2d} 스트라이프 정리 + 제목 맞춤")
            continue

        kept = []
        if spec.get("keep_pics"):
            kept = [sh for sh in slide.shapes if sh.__class__.__name__ == "Picture"]
            for sh in list(slide.shapes)[HEADER_N:]:
                if sh not in kept:
                    sh._element.getparent().remove(sh._element)
            for sh in kept:                      # 캡션 자리 확보
                if sh.top / 914400 + sh.height / 914400 > CAP_Y - 0.12:
                    sh.height = int((CAP_Y - 0.16 - sh.top / 914400) * 914400)
        elif i == COVER:
            for sh in list(slide.shapes):
                yy = sh.top / 914400
                if COVER_BAND[0] <= yy <= COVER_BAND[1]:
                    sh._element.getparent().remove(sh._element)
        else:
            clear_zone(slide)

        # 머리말 기하 — 2026-01-13 덱 배치에 맞춘다 (제목 전폭 · 정의는 구분선 아래 우측)
        if len(slide.shapes) <= 13:
            print(f"      ⚠ 머리말 도형이 {len(slide.shapes)}개뿐 — 기하 조정 건너뜀")
            continue
        slide.shapes[1].left, slide.shapes[1].width = Inches(0.57), Inches(8.90)
        slide.shapes[1].top, slide.shapes[1].height = Inches(0.26), Inches(0.50)
        slide.shapes[2].left, slide.shapes[2].top = Inches(5.02), Inches(0.94)
        slide.shapes[2].width, slide.shapes[2].height = Inches(4.46), Inches(0.42)
        slide.shapes[10].left, slide.shapes[10].top = Inches(5.78), Inches(1.40)
        slide.shapes[10].width = Inches(3.70)
        if spec.get("gloss"):
            set_text(slide.shapes[2], spec["gloss"])
        if spec.get("src"):
            set_text(slide.shapes[10], spec["src"])
        if spec.get("title"):
            set_text(slide.shapes[1], spec["title"])
        if spec.get("header"):
            set_text(slide.shapes[5], spec["header"])
        for k, b in enumerate(spec.get("bullets", [])):
            set_text(slide.shapes[7 + k * 2], b)

        if spec.get("extra_pics"):               # 기존 그림 옆에 추가
            n = len(kept) + len(spec["extra_pics"])
            gap = 0.22
            cw = (ZONE["w"] - gap * (n - 1)) / n
            for j, sh in enumerate(kept):
                bx = ZONE["x"] + j * (cw + gap)
                x, y, w, h = fit(sh.width, sh.height, bx, ZONE["y"], cw, ZONE["h"])
                sh.left, sh.top, sh.width, sh.height = (Inches(x), Inches(y), Inches(w), Inches(h))
            for j, rel in enumerate(spec["extra_pics"], start=len(kept)):
                p = FIG + rel
                with Image.open(p) as im:
                    iw, ih = im.size
                bx = ZONE["x"] + j * (cw + gap)
                x, y, w, h = fit(iw, ih, bx, ZONE["y"], cw, ZONE["h"])
                slide.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
        if spec.get("pics"):
            add_pics(slide, spec["pics"], zone=spec.get("pic_zone", ZONE))
        if spec.get("roster"):
            add_roster(slide)
        if spec.get("roster_boxes"):
            add_roster_boxes(slide)
        if spec.get("table"):
            t = spec["table"]
            add_table(slide, t["rows"], t.get("x", 0.62), t.get("y", 2.55),
                      t.get("w", 8.76), size=t.get("size", 11.5),
                      first_col_w=t.get("first_col_w"))
        if spec.get("fig_label"):
            labs = spec["fig_label"]
            labs = [labs] if isinstance(labs, str) else labs
            pics = [sh for sh in slide.shapes if sh.__class__.__name__ == "Picture"]
            for lab, sh in zip(labs, pics):
                add_fig_label(slide, lab, sh.left / 914400,
                              sh.top / 914400 + sh.height / 914400 + 0.02,
                              sh.width / 914400)
            notes = spec.get("fig_note")
            if notes:
                notes = [notes] if isinstance(notes, str) else notes
                for nt, sh in zip(notes, pics):
                    add_fig_note(slide, nt, sh.left / 914400,
                                 sh.top / 914400 + sh.height / 914400 + 0.30,
                                 sh.width / 914400)
        if spec.get("lines"):  # None 이면 건너뛴다
            add_lines(slide, spec["lines"], y=spec.get("lines_y", 2.75),
                      size=spec.get("lines_size", 13), lead=spec.get("lines_lead", 0.34))
        if spec.get("statement"):
            add_statement(slide, spec["statement"], y=spec.get("stmt_y", 3.6),
                          size=spec.get("stmt_size", 17.5),
                          x=spec.get("stmt_x", 1.05), w=spec.get("stmt_w", 7.90))
        if spec.get("cap"):
            add_caption(slide, spec["cap"])
        for j, floor in (() if i == COVER else ((1, 15.0), (5, 12.75))):   # 제목·구역 헤더는 한 줄 유지
            r = autofit_line(slide.shapes[j], floor_pt=floor)
            if r:
                print(f"      · shape[{j}] {r[0]:.2f}→{r[1]:.2f} pt" + ("" if r[2] else "  ⚠ 여전히 넘침"))
        print(f"  P{i:2d} 재구성 (pics={len(spec.get('pics',[]))+len(spec.get('extra_pics',[]))+len(kept)})")

    # 페이지 번호 재부여 — 재정렬하면 원본 번호가 그대로 따라온다
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and sh.top / 914400 > 6.9 and sh.left / 914400 > 8.5:
                set_text(sh, str(i))
                break
    prs.save(dst)
    print(f"\n저장: {dst}")


def selftest():
    ok = fail = 0
    def chk(name, cond):
        nonlocal ok, fail
        if cond: ok += 1
        else:
            fail += 1; print(f"  ✗ {name}")
    # 양성 — 비율 유지 맞춤
    x, y, w, h = fit(200, 100, 0, 0, 8, 4)
    chk("가로가 꽉 차는 경우", abs(w - 8) < 1e-9 and abs(h - 4) < 1e-9)
    x, y, w, h = fit(100, 200, 0, 0, 8, 4)
    chk("세로가 꽉 차는 경우", abs(h - 4) < 1e-9 and abs(w - 2) < 1e-9)
    chk("세로 맞춤 시 가로 중앙", abs(x - 3.0) < 1e-9)
    # 음성 ① — 박스를 넘치면 안 된다
    for pw, ph in ((3000, 40), (40, 3000), (1, 1), (1920, 1080)):
        x, y, w, h = fit(pw, ph, 0.5, 2.5, 8.9, 3.86)
        chk(f"음성: {pw}x{ph} 박스 이탈 없음",
            w <= 8.9 + 1e-9 and h <= 3.86 + 1e-9 and x >= 0.5 - 1e-9 and y >= 2.5 - 1e-9)
    # 음성 ② — 머리말을 건드리면 안 된다
    chk("음성: HEADER_N 이 푸터(13)보다 크다", HEADER_N > 13)
    # 음성 ③ — 캡션이 푸터(7.11)를 침범하면 안 된다
    chk("음성: 캡션이 푸터 위", CAP_Y + 0.44 < 7.11)
    # 음성 ④ — 내용 영역이 구분선(0.86)/캡션과 겹치면 안 된다
    chk("음성: 영역이 머리말 아래", ZONE["y"] > 2.2)
    chk("음성: 영역이 캡션 위에서 끝난다", ZONE["y"] + ZONE["h"] <= CAP_Y + 1e-9)
    # 음성 ⑤ — SPEC 이 가리키는 그림이 실제로 있어야 한다
    miss = [p for s in SPEC.values() for p in (s.get("pics", []) + s.get("extra_pics", []))
            if not os.path.exists(FIG + p)]
    chk(f"음성: 없는 그림 참조 0건 ({miss[:2]})", not miss)
    # 음성 ⑥ — KEEP_ZONE 과 SPEC 이 겹치면 안 된다 (지웠다 살리는 모순)
    bad_keep = [n for n in KEEP_ZONE
                if any(SPEC.get(n, {}).get(k) for k in ("pics", "extra_pics", "lines", "statement", "roster"))]
    chk(f"음성: KEEP_ZONE 장에 내용 배치 지시 없음 ({bad_keep})", not bad_keep)
    # 음성 ⑨ — 표지 띠가 푸터(7.09)·부제(5.57)를 건드리면 안 된다
    chk("음성: 표지 띠가 부제 위에서 끝난다", COVER_BAND[1] < 5.5)
    chk("음성: 표지 띠가 푸터를 안 건드린다", COVER_BAND[1] < 7.0)
    chk("음성: 표지에는 그림 지시가 없다", not SPEC.get(COVER, {}).get("pics"))
    # 음성 ⑪ — 표가 캡션·푸터를 침범하면 안 된다
    bad_t = [(n, round(sp_["table"].get("y", 2.55)
                       + max(0.225, 0.0235 * sp_["table"].get("size", 11.5))
                       * len(sp_["table"]["rows"]) + 0.035, 2))
             for n, sp_ in SPEC.items() if sp_.get("table")
             if sp_["table"].get("y", 2.55)
                + max(0.225, 0.0235 * sp_["table"].get("size", 11.5))
                * len(sp_["table"]["rows"]) + 0.035 > CAP_Y]
    chk(f"음성: 표가 캡션 줄을 넘지 않음 ({bad_t})", not bad_t)
    # 음성 ⑫ — 그림 라벨 수가 그림 수보다 많으면 안 된다
    bad_l = [n for n, sp_ in SPEC.items()
             if sp_.get("fig_label") and not isinstance(sp_["fig_label"], str)
             and len(sp_["fig_label"]) > len(sp_.get("pics", [])) + len(sp_.get("extra_pics", []))
             and not sp_.get("keep_pics")]
    chk(f"음성: 그림보다 라벨이 많은 장 없음 ({bad_l})", not bad_l)
    # 음성 ⑬ — 문장 상자가 슬라이드(10 in)를 벗어나면 안 된다
    bad_s = [(n, sp_.get("stmt_x", 1.05) + sp_.get("stmt_w", 7.90)) for n, sp_ in SPEC.items()
             if sp_.get("statement") and sp_.get("stmt_x", 1.05) + sp_.get("stmt_w", 7.90) > 9.55]
    chk(f"음성: 문장 상자 슬라이드 이탈 0건 ({bad_s})", not bad_s)
    # 음성 ⑩ — 용어 띠·출처 라벨이 너무 길면 머리말을 덮는다 (상자 5.26 in · 8.6 pt)
    gl = [(n, len(_MARK.sub(r"\2", sp_["gloss"]))) for n, sp_ in SPEC.items()
          if sp_.get("gloss") and len(_MARK.sub(r"\2", sp_["gloss"])) > 190]
    chk(f"음성: 용어 띠 길이 초과 0건 ({gl})", not gl)
    sc = [(n, len(sp_["src"])) for n, sp_ in SPEC.items() if sp_.get("src") and len(sp_["src"]) > 76]
    chk(f"음성: 출처 라벨 길이 초과 0건 ({sc})", not sc)
    # 음성 ⑦ — 제목/헤더/불릿이 한 줄 상자를 넘치면 안 된다 (2026-08-16 mock 렌더 실측 한도)
    LIM = dict(title=52, header=72)
    _plain = lambda t: _MARK.sub(r"\2", t)
    over = [(n, k, len(_plain(sp_[k]))) for n, sp_ in SPEC.items() for k, lim in LIM.items()
            if sp_.get(k) and len(_plain(sp_[k])) > lim]
    chk(f"음성: 제목·헤더 길이 초과 0건 ({over[:2]})", not over)
    ob = [(n, len(_plain(b))) for n, sp_ in SPEC.items()
          for b in sp_.get("bullets", []) if len(_plain(b)) > 95]
    chk(f"음성: 불릿 길이 초과 0건 ({ob[:2]})", not ob)
    # 음성 ⑧ — 폭 측정이 단조인가 (크기를 키우면 폭도 커져야 축소 루프가 종료된다)
    chk("음성: 폭이 pt 에 단조 증가", _text_pt("abc", 24) > _text_pt("abc", 12) > 0)
    chk("음성: 긴 문자열이 더 넓다", _text_pt("a" * 40, 18) > _text_pt("a" * 10, 18))
    chk("음성: 빈 문자열 폭 0", _text_pt("", 18) == 0)
    print(f"\nselftest: {ok} passed, {fail} failed")
    return 1 if fail else 0


if __name__ == "__main__":
    if "--selftest" in sys.argv:
        raise SystemExit(selftest())
    raise SystemExit(main() or 0)
