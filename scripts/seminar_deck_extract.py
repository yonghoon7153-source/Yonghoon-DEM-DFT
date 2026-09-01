#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""세미나 덱(.pptx) → 구조화 JSON.  웹앱 `/seminar` 이 읽는 정본을 만든다.

★ 왜 JSON 을 리포에 두는가: 웹앱이 요청마다 pptx 를 파싱하면 (a) python-pptx 를
  런타임 의존으로 만들고 (b) 발표 중 페이지가 파싱 오류로 죽을 수 있다.  덱이 바뀔 때만
  이 스크립트를 돌려 JSON 을 갱신하면 웹앱은 파일 하나만 읽으면 된다 — 기존 세미나
  페이지의 원칙("파일이 정본, 웹앱은 뷰어")을 그대로 따른다.

★ 분류는 **글꼴 크기**로 한다 (좌표가 아니라).  이 덱의 규약:
    26 pt bold  → kicker      (섹션 라벨: STEP 1 / Validation / SDCP …)
    15 pt bold  → title       (슬라이드 한 줄 주장)
    11.5 pt     → lead        (본문 문단)
    9–10 pt     → footer      (HANYANG / Battery Materials Lab. / 쪽번호) → 버린다
  그 밖의 크기는 **버리지 않고** `extra` 로 남긴다 — 조용한 결손을 만들지 않는다
  (덱 레이아웃이 바뀌면 extra 가 늘어나는 것으로 드러난다).

★ 발표 노트에서 `[Sources]` 블록을 떼어 `sources` 로 만든다.  대본과 출처를 한 화면에
  같이 놓는 것이 이 페이지의 목적이라, 노트 원문에서 분리해 둔다.

  python3 scripts/seminar_deck_extract.py                    # 기본 덱 → 기본 JSON
  python3 scripts/seminar_deck_extract.py --deck a.pptx --out b.json
  python3 scripts/seminar_deck_extract.py --selftest
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys

DECK_DEFAULT = os.path.join('docs', 'seminar', 'seminar_20260806_DEM_MPM_SDCP_section.pptx')
OUT_DEFAULT = os.path.join('docs', 'seminar', 'seminar_deck.json')

#: 글꼴 크기(pt) → 역할.  덱 규약이 바뀌면 여기만 고친다.
SIZE_ROLE = {26.0: 'kicker', 15.0: 'title', 11.5: 'lead'}

#: 푸터 보일러플레이트 — 화면에 다시 그릴 이유가 없다.
FOOTER_RE = re.compile(r'^(HANYANG UNIVERSITY|Battery Materials Lab\.?|\d{1,2})\s*$', re.I)

#: 그림 자리표시자 — 슬라이드에 무엇이 들어갈 자리인지 화면에 남긴다.
FIGURE_RE = re.compile(r'^\s*FIGURE SLOT\b', re.I)


def _para_size(par):
    """문단의 대표 글꼴 크기(pt).  런마다 다르면 **최댓값** (제목이 잘리지 않게)."""
    sizes = [r.font.size.pt for r in par.runs if r.font.size is not None]
    return max(sizes) if sizes else None


def split_sources(notes):
    """노트 → (본문, [출처…]).  `[Sources]` 이후 줄의 `- ` 항목을 걷는다."""
    if not notes:
        return '', []
    m = re.search(r'^\s*\[Sources\]\s*$', notes, re.M)
    if not m:
        return notes.strip(), []
    body, tail = notes[:m.start()], notes[m.end():]
    srcs = [ln.strip().lstrip('-').strip() for ln in tail.splitlines() if ln.strip().startswith('-')]
    return body.strip(), [s for s in srcs if s]


def extract_slide(slide, number):
    """한 슬라이드 → dict.  분류 못 한 텍스트는 버리지 않고 extra 로 남긴다."""
    out = {'n': number, 'kicker': '', 'title': '', 'lead': [],
           'tables': [], 'charts': [], 'figure_slots': [], 'extra': [],
           'notes': '', 'sources': []}
    for sh in slide.shapes:
        # ★ 차트를 버리지 않는다 — slide 13 의 σ_SDCP 민감도가 차트로만 들어 있어,
        #   빠뜨리면 이 발표의 메커니즘 근거가 화면에서 통째로 사라진다.
        if getattr(sh, 'has_chart', False) and sh.has_chart:
            ch = sh.chart
            try:
                cats = [str(c) for c in ch.plots[0].categories]
            except Exception:                                  # noqa: BLE001
                cats = []
            series = []
            for sr in ch.series:
                series.append({'name': getattr(sr, 'name', '') or '',
                               'values': [None if v is None else float(v) for v in sr.values]})
            out['charts'].append({'title': (ch.chart_title.text_frame.text
                                            if ch.has_title else ''),
                                  'categories': cats, 'series': series})
            continue
        if getattr(sh, 'has_table', False) and sh.has_table:
            rows = [[(c.text or '').strip() for c in r.cells] for r in sh.table.rows]
            rows = [r for r in rows if any(r)]
            if rows:
                out['tables'].append(rows)
            continue
        if not sh.has_text_frame:
            continue
        for par in sh.text_frame.paragraphs:
            t = (par.text or '').strip()
            if not t or FOOTER_RE.match(t):
                continue
            if FIGURE_RE.match(t):
                out['figure_slots'].append(t)
                continue
            role = SIZE_ROLE.get(_para_size(par))
            if role == 'kicker' and not out['kicker']:
                out['kicker'] = t
            elif role == 'title' and not out['title']:
                out['title'] = t
            elif role == 'lead':
                out['lead'].append(t)
            else:
                out['extra'].append(t)
    notes = ''
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
        notes = slide.notes_slide.notes_text_frame.text or ''
    out['notes'], out['sources'] = split_sources(notes)
    return out


def extract(deck_path):
    from pptx import Presentation                       # 이 스크립트에서만 필요
    pres = Presentation(deck_path)
    slides = [extract_slide(s, i + 1) for i, s in enumerate(pres.slides)]
    return {
        'deck': os.path.relpath(deck_path).replace(os.sep, '/'),
        'n_slides': len(slides),
        'slides': slides,
        # 화면이 "무엇이 안 잡혔는지" 를 바로 보여줄 수 있게 집계해 둔다.
        'unclassified': sum(len(s['extra']) for s in slides),
        'n_charts': sum(len(s['charts']) for s in slides),
        'no_notes': [s['n'] for s in slides if not s['notes']],
    }


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ap.add_argument('--deck', default=os.path.join(here, DECK_DEFAULT))
    ap.add_argument('--out', default=os.path.join(here, OUT_DEFAULT))
    ap.add_argument('--selftest', action='store_true')
    a = ap.parse_args(argv)
    if a.selftest:
        return _selftest()
    if not os.path.exists(a.deck):
        sys.exit(f'덱 없음: {a.deck}')
    data = extract(a.deck)
    #  ★ 기존 산출물의 `_RETRACTED` 는 **이어받는다** (2026-09-01).  이 키는 덱에서
    #    뽑히는 것이 아니라 사람이 붙인 철회 표지라, 추출기를 다시 돌리면 조용히
    #    사라진다.  실사고: `56dcbca1` 이 pptx 1번 장에 배너를 끼웠는데 JSON 을 다시
    #    안 뽑아 **화면은 배너 없는 14장짜리 옛 덱**을 계속 보여주고 있었다.
    if os.path.exists(a.out):
        try:
            with open(a.out, encoding='utf-8') as f:
                _old = json.load(f)
            if _old.get('_RETRACTED') and not data.get('_RETRACTED'):
                data = {'_RETRACTED': _old['_RETRACTED'], **data}
                print('  ℹ 이전 산출물의 _RETRACTED 표지를 이어받았다')
        except (OSError, ValueError) as _e:
            print(f'  ⚠ 이전 산출물을 못 읽었다 ({type(_e).__name__}) — 표지 이어받기 생략')
    os.makedirs(os.path.dirname(a.out), exist_ok=True)
    with open(a.out, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=1)
    print(f"{data['n_slides']} 슬라이드 → {a.out}")
    for s in data['slides']:
        print(f"  {s['n']:>2}. [{s['kicker'] or '—'}] {s['title'][:58]}"
              f"   표{len(s['tables'])} 노트{len(s['notes'])}자 출처{len(s['sources'])}")
    if data['unclassified']:
        print(f"\n⚠ 분류 안 된 텍스트 {data['unclassified']} 조각 (extra 로 보존 — 화면에도 나온다)")
    if data['no_notes']:
        print(f"⚠ 노트 없는 슬라이드: {data['no_notes']}")
    return 0


def _selftest():
    n = [0, 0]

    def ok(name, cond):
        n[1] += 1
        n[0] += bool(cond)
        print(f'  {"PASS" if cond else "FAIL"}  {name}')

    body, srcs = split_sources("설명 문장.\n\n[Sources]\n- docs/a.md\n- scripts/b.py\n")
    ok('1) [Sources] 를 노트 본문에서 떼어낸다',
       body == '설명 문장.' and srcs == ['docs/a.md', 'scripts/b.py'])
    ok('2) [Sources] 가 없으면 노트 전체가 본문',
       split_sources('그냥 대본') == ('그냥 대본', []))
    ok('3) 빈 노트도 안전', split_sources('') == ('', []) and split_sources(None) == ('', []))
    ok('4) 푸터 보일러플레이트를 거른다',
       all(FOOTER_RE.match(x) for x in ('HANYANG UNIVERSITY', 'Battery Materials Lab.', '13'))
       and not FOOTER_RE.match('STEP 1'))
    ok('5) ★ 쪽번호만 거르고 숫자로 시작하는 문장은 남긴다',
       FOOTER_RE.match('7') and not FOOTER_RE.match('7 개 채널을 푼다'))
    ok('6) FIGURE SLOT 인식', FIGURE_RE.match('FIGURE SLOT — webapp 3D bed')
       and not FIGURE_RE.match('figure 후보만 표시'))
    ok('7) 크기→역할 표가 덱 규약과 같다',
       SIZE_ROLE == {26.0: 'kicker', 15.0: 'title', 11.5: 'lead'})

    here = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    deck = os.path.join(here, DECK_DEFAULT)
    if os.path.exists(deck):
        try:
            d = extract(deck)
            ok('8) 실제 덱 추출 — 슬라이드 수 > 0', d['n_slides'] > 0)
            #  ★ 옛 판은 `s['n'] == 1` 로 표지를 면제했다.  그런데 `56dcbca1` 이 철회
            #    배너를 **1번 장으로 앞에 끼우자** 표지가 2번이 되어 이 검사가 빨간불이
            #    됐고, 아무 레인에도 없어 아무도 몰랐다.  ⇒ 면제 기준을 **위치가 아니라
            #    내용**으로 바꾼다: 본문(lead·표·차트)이 없는 장 하나까지만 봐준다.
            _untitled = [s for s in d['slides'] if not s['title']]
            ok('9) ★ 제목 없는 장은 본문 없는 표지 하나뿐이다 (크기-분류가 통했다)',
               len(_untitled) <= 1
               and all(not (s['lead'] or s['tables'] or s['charts']) for s in _untitled))
            ok('10) ★ 노트가 있는 슬라이드가 대다수 (대본이 이 페이지의 핵심)',
               sum(1 for s in d['slides'] if s['notes']) >= d['n_slides'] - 1)
            ok('11) ★ 차트를 버리지 않는다 (σ_SDCP 민감도가 차트로만 들어 있다)',
               d['n_charts'] >= 1
               and any(s['charts'] and s['charts'][0]['series'] for s in d['slides']))
        except ImportError:
            print('  SKIP 8~10) python-pptx 없음')
    else:
        print('  SKIP 8~10) 덱 파일 없음')
    print(f'\nseminar_deck_extract selftest: {n[0]}/{n[1]} PASS')
    return 0 if n[0] == n[1] else 1


if __name__ == '__main__':
    raise SystemExit(main())
