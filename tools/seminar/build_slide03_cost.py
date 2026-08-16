import sys
sys.path.insert(0,'/home/user/Yonghoon-DEM-DFT/tools/seminar')
import rebuild_cascade_deck as M
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(0x00,0x70,0xC0); INK = RGBColor(0x1f,0x29,0x37); MUT = RGBColor(0x6b,0x72,0x80)

def dashed_box(sl, title, x, y, w, h):
    """6/15 덱 관례 — 파란 점선 둥근 사각형 + 테두리를 끊고 앉는 흰 배경 제목."""
    b = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.adjustments[0] = 0.05
    b.fill.background(); b.line.color.rgb = BLUE; b.line.width = Pt(1.75)
    b.line.dash_style = 4; b.shadow.inherit = False
    t = sl.shapes.add_textbox(Inches(x+0.55), Inches(y-0.155), Inches(3.0), Inches(0.31))
    t.fill.solid(); t.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF); t.line.fill.background()
    pr = t.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    r = pr.add_run(); r.text = title
    r.font.size, r.font.bold, r.font.name, r.font.color.rgb = Pt(13.5), True, "Arial", BLUE

def step(sl, n, label, note, y, xn=1.15, xl=1.52, xa=5.25):
    for x, w, txt, sz, col in ((xn,0.34,f"{n}.",14,INK), (xl,3.60,label,14,INK), (xa,3.60,f"({note})",13,MUT)):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = txt
        r.font.size, r.font.name, r.font.color.rgb = Pt(sz), "Arial", col

prs = Presentation('slide3_only.pptx'); sl = prs.slides[0]
M.clear_zone(sl)
sl.shapes[1].left, sl.shapes[1].width   = Inches(0.57), Inches(8.90)
sl.shapes[1].top,  sl.shapes[1].height  = Inches(0.26), Inches(0.50)
sl.shapes[2].left, sl.shapes[2].top     = Inches(0.60), Inches(6.86)
sl.shapes[2].width, sl.shapes[2].height = Inches(8.86), Inches(0.24)
sl.shapes[10].left, sl.shapes[10].top   = Inches(5.78), Inches(6.62)
sl.shapes[10].width = Inches(3.68)

M.set_text(sl.shapes[1], "Cost of the standard doping route")
M.set_text(sl.shapes[5], "Motivation (2): The route we already use is accurate but slow")
M.set_text(sl.shapes[7], "Several calculations were run for the project, and each of them is accurate.")
M.set_text(sl.shapes[9], "The cost per composition is [r]large[/r], and it grows with every candidate added.")
M.set_text(sl.shapes[2], "MLIP: machine-learned interatomic potential  ·  EOS: energy as a function of cell volume  ·  "
                         "Multi-seed: the same run repeated from different starting arrangements")
M.set_text(sl.shapes[10], "Our standard route, per doped composition")

dashed_box(sl, "Standard route per composition", 0.62, 2.66, 8.76, 2.30)
for i,(lab,note) in enumerate([("Enumerate substitutions","site × charge recipe"),
                               ("MLIP screen + anneal","Li ordering"),
                               ("DFT validation","EOS, elastic, band structure"),
                               ("Long MD","conductivity, multi-seed")]):
    step(sl, i+1, lab, note, 2.96 + i*0.50)

M.add_statement(sl, "So a low-cost screening route was devised, and run over 91 candidates.",
                y=5.42, size=17.5)
for j, floor in ((1,15.0),(5,12.75)): M.autofit_line(sl.shapes[j], floor_pt=floor)
prs.save('Slide03_cost.pptx'); print('saved')
