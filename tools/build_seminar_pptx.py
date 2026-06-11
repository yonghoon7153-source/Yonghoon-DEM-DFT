#!/usr/bin/env python3
"""
Build paper #1 seminar PPT from content in
kb/papers/lpscl_vs_lpscl16_seminar_v1.md

Generates: kb/papers/lpscl_vs_lpscl16_seminar_v1.pptx
21 slides, 16:9, Korean (한국어 발표용)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# Colors
HANYANG_BLUE = RGBColor(0x00, 0x38, 0x76)
HIGHLIGHT_RED = RGBColor(0xC8, 0x10, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
LIGHT_BLUE = RGBColor(0xE6, 0xEE, 0xF7)
LIGHT_RED = RGBColor(0xFB, 0xE4, 0xE7)
LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xE6)
LIGHT_ORANGE = RGBColor(0xFD, 0xF0, 0xE0)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, font="맑은 고딕"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_title(slide, text, color=HANYANG_BLUE, size=24):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             text, size=size, bold=True, color=color)
    # bottom line under title
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Emu(20000),
             fill=color)


def add_footer(slide, num, total=26):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.3),
             "LPSCl ↔ LPSCl₁.₆ — Paper #1 Preview · Yonghoon Kim · 2026.06",
             size=9, color=GRAY)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.3),
             f"{num} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, data, first_row_header=True,
              first_col_header=False, header_color=HANYANG_BLUE,
              header_text=RGBColor(0xFF, 0xFF, 0xFF), font_size=11,
              row_colors=None, bold_first_col=False, col_widths=None):
    """data = list of rows; each row = list of strings.
    row_colors: dict {row_idx: RGBColor} for highlight rows."""
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths is not None:
        total = sum(col_widths)
        for c, frac in enumerate(col_widths):
            tbl.columns[c].width = Emu(int(width * frac / total))
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.06)
            tf.margin_top = tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.name = "맑은 고딕"
            if r == 0 and first_row_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.bold = True
                run.font.color.rgb = header_text
            elif row_colors and r in row_colors:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_colors[r]
                run.font.color.rgb = BLACK
                if c == cols - 1 or "★" in str(val):
                    run.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.color.rgb = BLACK
            if bold_first_col and c == 0 and r > 0:
                run.font.bold = True
    return tbl


def slide_blank():
    return prs.slides.add_slide(BLANK)


# =====================================================
# SLIDE 1 — Title
# =====================================================
s = slide_blank()
# Big title
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.2),
         "LPSCl₁.₆가 빠르고 단단한 이유:\n— 전자구조가 아니라 무질서다 —",
         size=36, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
# Subtitle EN
add_text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.8),
         "A multi-probe DFT/AIMD comparison of\nLi₆PS₅Cl vs Li₅.₄PS₄.₄Cl₁.₆",
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
# Thesis box
add_rect(s, Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.4),
         fill=LIGHT_GRAY)
# vertical accent
add_rect(s, Inches(2.0), Inches(4.2), Emu(40000), Inches(1.4),
         fill=HANYANG_BLUE)
add_text(s, Inches(2.2), Inches(4.35), Inches(8.9), Inches(0.4),
         "Thesis", size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(2.2), Inches(4.75), Inches(8.9), Inches(0.9),
         '차이는 "전자구조"가 아니라\nLi 공공 + 4d-Cl anti-site 무질서에서 온다.',
         size=18, bold=True, color=BLACK)
# Author footer
add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Emu(15000), fill=GRAY)
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
         "안용훈 · BML Lab (한양대) · 지도 김광범 교수\n2026년 6월 · Internal Seminar Preview",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 2 — Scope: 두 시스템
# =====================================================
s = slide_blank()
add_title(s, "2. 비교 대상: 두 argyrodite")

# Left panel (comp1)
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.2),
         fill=LIGHT_BLUE, line=HANYANG_BLUE, line_w=1.5)
add_text(s, Inches(0.95), Inches(1.55), Inches(5.5), Inches(0.5),
         "LPSCl (comp1)", size=22, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.95), Inches(2.05), Inches(5.5), Inches(0.45),
         "Li₆PS₅Cl", size=18, color=BLACK)
add_text(s, Inches(0.95), Inches(2.6), Inches(5.5), Inches(0.4),
         "cubic F-43m · 4 f.u. · 52 atoms", size=13, color=GRAY)
add_text(s, Inches(0.95), Inches(3.3), Inches(5.5), Inches(2.8),
         "• ordered Li (vacancy 없음)\n"
         "• Cl 전부 4a 자리\n"
         "• free S²⁻ 4d 자리\n"
         "• 단일 환경 (모든 P, Cl equivalent)",
         size=14)
add_text(s, Inches(0.95), Inches(5.9), Inches(5.5), Inches(0.5),
         "V₀ = 1016.62 Å³", size=12, color=GRAY)

# Right panel (modelc)
add_rect(s, Inches(6.75), Inches(1.4), Inches(5.8), Inches(5.2),
         fill=LIGHT_RED, line=HIGHLIGHT_RED, line_w=1.5)
add_text(s, Inches(6.9), Inches(1.55), Inches(5.5), Inches(0.5),
         "LPSCl₁.₆ (modelc)", size=22, bold=True, color=HIGHLIGHT_RED)
add_text(s, Inches(6.9), Inches(2.05), Inches(5.5), Inches(0.45),
         "Li₅.₄PS₄.₄Cl₁.₆", size=18, color=BLACK)
add_text(s, Inches(6.9), Inches(2.6), Inches(5.5), Inches(0.4),
         "rhombohedral R3m · 5 f.u. · 62 atoms", size=13, color=GRAY)
add_text(s, Inches(6.9), Inches(3.3), Inches(5.5), Inches(2.8),
         "• Li 공공 0.6 / f.u. (vacancy)\n"
         "• Cl 4a + 4d anti-site (AS)\n"
         "• 4d 자리: S²⁻ + Cl 혼합\n"
         "• 두 disorder source 공존",
         size=14)
add_text(s, Inches(6.9), Inches(5.9), Inches(5.5), Inches(0.5),
         "V₀ = 1216.44 Å³", size=12, color=GRAY)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "→ 동일 protocol · 동일 §8 multi-probe로 paired 비교",
         size=14, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 2)


# =====================================================
# SLIDE 3 — Pipeline 3-tier
# =====================================================
s = slide_blank()
add_title(s, "3. Pipeline: MLIP screen → DFT confirm → multi-probe")

# 3 boxes
tiers = [
    ("Tier 1 · MLIP screening", "hours", LIGHT_BLUE,
     ["• Halogen enumerate (45 configs)",
      "• Li sublattice screen (top-5 × 20)",
      "• 500 K Langevin anneal → champion",
      "• MLIP EOS pre-scan → V₀ 범위",
      "  UMA-s-1p1 (omat)"]),
    ("Tier 2 · DFT validation", "days", RGBColor(0xCD, 0xDC, 0xEE),
     ["• BM3 EOS 11 volumes (V/V₀ 0.96–1.06)",
      "• V₀ confirmation BFGS, force <5e-3",
      "• k-mesh 수렴 보장 (k×L ≥ 40 Å)",
      "• V₀, B₀ paper-grade (<1 GPa)",
      "  PBE + USPP, ecut 60 / 480 Ry"]),
    ("Tier 3 · §8 multi-probe", "weeks", RGBColor(0xA5, 0xC4, 0xE4),
     ["• structure  bonds · Voronoi · BVSE",
      "• electronic DOS · bands · ELF",
      "• bonding    Bader · LOBSTER ICOHP",
      "• transport  AIMD 600/800/1000 K",
      "• mechanical stress-strain Cij"]),
]
ty = 1.35
for title, time, color, bullets in tiers:
    add_rect(s, Inches(0.8), Inches(ty), Inches(11.7), Inches(1.65),
             fill=color, line=HANYANG_BLUE, line_w=0.5)
    add_text(s, Inches(1.0), Inches(ty + 0.1), Inches(8.5), Inches(0.4),
             title, size=16, bold=True, color=HANYANG_BLUE)
    add_text(s, Inches(10.5), Inches(ty + 0.1), Inches(2.0), Inches(0.4),
             time, size=12, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(1.0), Inches(ty + 0.5), Inches(11.5), Inches(1.1),
             "\n".join(bullets), size=12)
    ty += 1.85

add_text(s, Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.3),
         "→ 4500 MLIP configs screened · 1 champion DFT-validated · 13 paper-grade probes",
         size=13, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 3)


# =====================================================
# SLIDE 4 — Headline 4-message table
# =====================================================
s = slide_blank()
add_title(s, "4. Headline — 한 페이지로 보는 두 시스템 (paper-grade)")

headline = [
    ["항목", "LPSCl", "LPSCl₁.₆", "메시지"],
    ["밴드갭 (eV)", "1.76", "1.82", "거의 동일  ★ M1"],
    ["AIMD Ea (eV) ⁴ᶠᵘ", "0.253", "0.224", "modelc 낮음 (Cl-rich Ea↓)"],
    ["D(600 K) cm²/s", "3.09e-6", "7.90e-6", "modelc 2.5× 빠름  ★ M2"],
    ["D₀ prefactor", "4.11e-4", "5.8e-4", "modelc 1.4× — vacancy carrier"],
    ["ICOHP Li–Cl (eV)", "−1.86", "−2.10", "+13% 강화  ★ M3"],
    ["Bader Li (e)", "+0.874", "+0.882", "거의 동일"],
    ["E_VRH relaxed-ion (GPa)", "22.06", "27.66", "+25.4% 단단  ★ M4"],
    ["B₀ (BM-EOS, GPa)", "26.23", "21.71", "hydrostatic 반대 (vacancy soft)"],
    ["Zener A", "1.14", "1.44", "비등방성 ↑"],
]
row_colors = {
    1: LIGHT_BLUE,            # M1
    3: LIGHT_GREEN, 4: LIGHT_GREEN,  # M2
    5: LIGHT_ORANGE,          # M3
    7: LIGHT_RED,             # M4
}
add_table(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.7),
          headline, row_colors=row_colors, font_size=12)

add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
         "→ 4 messages = 전자구조 둔감 · σ는 dual mechanism · ionic glue 강화 · vacancy paradox 해소",
         size=13, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3),
         "⁴ᶠᵘ comp1 4 f.u. natural cubic (2026-06-11 갱신) — Schlem 2020 LPSCl 0.25 eV EXACT 매칭",
         size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 4)


# =====================================================
# SLIDE 5 — M1 전자구조 둔감
# =====================================================
s = slide_blank()
add_title(s, "5. M1: 전자구조는 두 시스템에서 거의 동일")

# Left: DOS placeholder
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.5),
         fill=LIGHT_GRAY, line=GRAY, line_w=0.5)
add_text(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(0.5),
         "[ DOS overlay plot ]\n(comp1 파랑 / modelc 빨강)\nE − EF 기준",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Right: table
m1_data = [
    ["항목", "LPSCl", "LPSCl₁.₆"],
    ["gap (eV)", "1.76", "1.82"],
    ["Δgap (eV)", "—", "+0.06"],
    ["EF (eV)", "2.82", "2.45"],
    ["VBM peak (eV)", "1.64", "1.84"],
    ["CBM peak (eV)", "5.44", "5.72"],
    ["VBM character", "S 3p 91% + Li p 6%", "S 3p 92% + Li p 6%"],
    ["CBM character", "S p + P s + Li p", "동일 패턴"],
]
add_table(s, Inches(6.7), Inches(1.4), Inches(6.2), Inches(3.6),
          m1_data, font_size=11)

# Bullets
add_text(s, Inches(6.7), Inches(5.15), Inches(6.2), Inches(0.9),
         "• 갭 차이 Δ = 0.06 eV (작음)\n"
         "• VBM/CBM 궤도 character 거의 동일\n"
         "• 작은 차이로는 σ 3× · E 25% 못 설명",
         size=11)

# Footnote
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
         "※ modelc EF(2.45) < VBM(2.72): 국소 defect-band 0.74 states\n"
         "  (comp1 동일 구간 0.037, 20× 적음) — vacancy + 4d-Cl AS의 전자적 흔적",
         size=10, color=GRAY)
add_footer(s, 5)


# =====================================================
# SLIDE 6 — M2 σ dual mechanism (v3 ACTIVE)
# =====================================================
s = slide_blank()
add_title(s, "6. M2: σ 차이 = barrier ↓ + prefactor ↑ 둘 다 작용")

# Left: Arrhenius plot placeholder
add_rect(s, Inches(0.5), Inches(1.4), Inches(5.5), Inches(4.5),
         fill=LIGHT_GRAY, line=GRAY, line_w=0.5)
add_text(s, Inches(0.5), Inches(3.4), Inches(5.5), Inches(0.5),
         "[ Arrhenius plot ]\nln D vs 1000/T\n(comp1 vs modelc, 3 points)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Right: table
m2_data = [
    ["T (K)", "LPSCl 4fu", "LPSCl₁.₆"],
    ["600", "3.09e-6", "7.90e-6"],
    ["800", "1.03e-5", "2.05e-5"],
    ["1000", "2.20e-5", "4.55e-5"],
    ["Ea (eV) ★", "0.253", "0.224"],
    ["D₀ (cm²/s)", "4.11e-4", "5.8e-4"],
    ["R²", "0.9998", "0.992"],
    ["Schlem 2020 exp.", "0.25 ✓", "0.22 ✓"],
]
row_colors_m2 = {4: LIGHT_GREEN, 5: LIGHT_GREEN, 7: LIGHT_BLUE}
add_table(s, Inches(6.3), Inches(1.4), Inches(6.7), Inches(4.0),
          m2_data, row_colors=row_colors_m2, font_size=11)

# 2-mechanism box
add_rect(s, Inches(0.5), Inches(6.0), Inches(12.5), Inches(1.0),
         fill=LIGHT_GREEN)
add_text(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.45),
         "σ(modelc) / σ(comp1) = 2.5× at 600 K  =  Ea 효과 1.75×  ×  D₀ 효과 1.41×",
         size=14, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.45),
         "→ Cl-rich가 lower Ea (Minafra direction) + higher D₀ (vacancy carrier) 둘 다 작용",
         size=12, color=BLACK, align=PP_ALIGN.CENTER)
add_footer(s, 6)


# =====================================================
# SLIDE 6a — Mechanism cartoon
# =====================================================
s = slide_blank()
add_title(s, "6a. 왜 carrier가 이기는가 — 두 효과의 경쟁")

# Two panels
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.0),
         fill=LIGHT_BLUE, line=HANYANG_BLUE, line_w=1)
add_text(s, Inches(0.95), Inches(1.55), Inches(5.2), Inches(0.5),
         "LPSCl (ordered)", size=18, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.95), Inches(2.1), Inches(5.2), Inches(2.7),
         "● ● ● ● ● ●  ← Li 가득 참\n\n"
         "Li hopping 위해서는\n"
         "옆자리 비기를 '대기'\n\n"
         "→ carrier 부족",
         size=14)

add_rect(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.0),
         fill=LIGHT_RED, line=HIGHLIGHT_RED, line_w=1)
add_text(s, Inches(7.15), Inches(1.55), Inches(5.2), Inches(0.5),
         "LPSCl₁.₆ (vacancy + AS)", size=18, bold=True, color=HIGHLIGHT_RED)
add_text(s, Inches(7.15), Inches(2.1), Inches(5.2), Inches(2.7),
         "● ● ○ ● ● ○  ← 공공 ○\n\n"
         "빈자리 항상 근처\n"
         "→ 즉시 hopping 가능\n\n"
         "→ ~8× 많은 path",
         size=14)

# Arrhenius formula
add_rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.0),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.4),
         "D = D₀ · exp(−Ea/kT)",
         size=18, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
         "↑ 1.4× (vacancy carrier)        ↑ 1.75× at 600K (Cl-rich barrier ↓)",
         size=12, color=BLACK, align=PP_ALIGN.CENTER)
add_footer(s, 7)


# =====================================================
# SLIDE 6b — Disorder ensemble 2×2
# =====================================================
s = slide_blank()
add_title(s, "6b. Disorder Ensemble — matched-d에서 Ea 동일")

d_data = [
    ["Ea (eV)", "LPSCl (4 f.u.)", "LPSCl₁.₆ (5 f.u.)"],
    ["clean (d=0)", "0.253 (★ Schlem)", "0.224 (★ Schlem)"],
    ["disordered (d≈0.4–0.5, n=3)", "0.177 ± 0.027", "0.173 ± 0.039"],
    ["", "← matched disorder Ea 동일 (Δ4 meV)", ""],
]
add_table(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2),
          d_data, font_size=12,
          row_colors={1: LIGHT_BLUE, 2: LIGHT_GREEN})

add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.5),
         "• 4 f.u. clean (Schlem 정확 매칭) → Cl-rich가 본질적으로 lower Ea\n"
         "• matched disorder에서 두 시스템 Ea 동일 (~0.18 eV)\n"
         "• 추가 disorder 주입 시 Ea가 더 떨어짐 (Minafra/Kraft mechanism 확인)\n"
         "• σ 3× 차이의 출처: D₀ (vacancy carrier) + Ea (disorder-induced 평탄화)",
         size=14)

add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
         "※ d=0.0 ordered LPSCl는 실험·우리 모두 disorder 일부 포함 (ball-milled),\n"
         "  pure ordered limit (1.17 eV)는 kinetically inaccessible artifact",
         size=10, color=GRAY)
add_footer(s, 8)


# =====================================================
# SLIDE 6c — No T_cross
# =====================================================
s = slide_blank()
add_title(s, "6c. 저온 특성 — modelc wins at ALL T, no trade-off")

trate_data = [
    ["T (K)", "Ea contrib.", "D₀ contrib.", "σ_ratio", "비고"],
    ["1000", "1.40×", "1.41×", "1.97×", "측정"],
    ["800", "1.51", "1.41", "2.12", "측정"],
    ["600", "1.75", "1.41", "2.47", "측정값 2.5× ✓"],
    ["500", "1.97", "1.41", "2.78", "외삽"],
    ["400", "2.36", "1.41", "3.33", "외삽"],
    ["300 (RT)", "3.07", "1.41", "4.33  ★", "외삽 (Zuo 2.4× cf)"],
    ["200", "4.94", "1.41", "6.96", "외삽 (불확실)"],
]
add_table(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(3.7),
          trate_data, font_size=11,
          row_colors={3: LIGHT_GREEN, 6: LIGHT_RED})

add_text(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5),
         "• T_cross 없음 — Ea와 D₀ 둘 다 modelc 우세 같은 방향\n"
         "• 저온일수록 σ ratio ↑ (Arrhenius)\n"
         "• Zuo 2023 RT 측정 2.4× ↔ 우리 외삽 4.3× — 자릿수 정합 ±30%\n"
         "• 'vacancy 양날' framing (v2 5fu 기반) 무효 — Cl-rich 저온도 우위",
         size=12)
add_footer(s, 9)


# =====================================================
# SLIDE 7 — M3 ionic glue
# =====================================================
s = slide_blank()
add_title(s, "7. M3: Li–anion ionic glue가 LPSCl₁.₆에서 강해진다")

m3_data = [
    ["결합", "LPSCl", "LPSCl₁.₆", "Δ%"],
    ["P–S (PS₄ covalent)", "−5.94", "−6.00", "+1% 불변"],
    ["Li–Cl  ★", "−1.86", "−2.10", "+13%"],
    ["Li–S  ★", "−1.59", "−1.72", "+8%"],
    ["S–S (cage)", "−0.11", "−0.11", "~0"],
]
add_table(s, Inches(0.6), Inches(1.4), Inches(6.3), Inches(2.3),
          m3_data, font_size=11, row_colors={2: LIGHT_ORANGE, 3: LIGHT_ORANGE})

# Per-bond decomposition
decomp = [
    ["Li–Cl 분해 단계", "per-bond", "범위"],
    ["comp1 baseline (전부 4a)", "−1.855", "Cl 전부"],
    ["modelc 4a (+vacancy)", "−2.026", "Cl 90% (+9.2%)"],
    ["modelc 4d AS (+anti-site)", "−2.836  ★", "Cl 10% (+40%)"],
]
add_table(s, Inches(7.1), Inches(1.4), Inches(5.7), Inches(2.3),
          decomp, font_size=10,
          row_colors={3: LIGHT_ORANGE})

# Contributions
add_rect(s, Inches(0.6), Inches(4.0), Inches(12.2), Inches(1.4),
         fill=LIGHT_ORANGE)
add_text(s, Inches(0.6), Inches(4.05), Inches(12.2), Inches(0.4),
         "+13% 평균 Li–Cl 강화 기여도",
         size=14, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.5), Inches(12.2), Inches(0.9),
         "• Vacancy field: 69% — Cl 90%에 균일 적용 (+9.2%/bond)\n"
         "• 4d-Cl Anti-site: 31% — Cl 10%에만 intense (+40%/bond)",
         size=13, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(1.0),
         "• Cl 치환은 결합 '강화' (약화 아님)\n"
         "• Bader · LOBSTER · Wilkening 세 독립 probe 합의",
         size=12)
add_footer(s, 10)


# =====================================================
# SLIDE 8 — M4 vacancy paradox
# =====================================================
s = slide_blank()
add_title(s, "8. M4: clamped 동일, relaxed-ion +25% — vacancy paradox 해소")

m4_main = [
    ["E_VRH (GPa)", "LPSCl", "LPSCl₁.₆", "Δ%", "실험 (Kim 2025)"],
    ["clamped-ion (frozen)", "52.31", "52.30", "0% — paradox", "—"],
    ["★ relaxed-ion (Born screening)", "22.06", "27.66", "+25%", "✓ ~23"],
]
add_table(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(1.5),
          m4_main, font_size=12, row_colors={2: LIGHT_RED})

decomp_m4 = [
    ["modulus", "LPSCl", "LPSCl₁.₆", "Δ%"],
    ["B_VRH", "25.5", "23.4", "−8% (vacancy soft)"],
    ["G_VRH  ★", "8.1", "10.6", "+30% (shear stiff)"],
    ["E_VRH", "22.06", "27.66", "+25%"],
    ["ν (Poisson)", "0.36", "0.30", "비등방화"],
    ["Zener A", "1.14", "1.44", "+26%"],
]
add_table(s, Inches(0.6), Inches(3.1), Inches(7.0), Inches(2.7),
          decomp_m4, font_size=10, row_colors={2: LIGHT_RED})

# Right side: causal chain
add_rect(s, Inches(7.8), Inches(3.1), Inches(5.0), Inches(2.7),
         fill=LIGHT_GRAY)
add_text(s, Inches(7.95), Inches(3.2), Inches(4.8), Inches(0.4),
         "M1 ↔ M3 ↔ M4 인과 chain",
         size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(7.95), Inches(3.6), Inches(4.8), Inches(2.1),
         "• PS₄ = rigid blocks (M1, 불변)\n"
         "• Li–anion = mortar (M3 +13%)\n"
         "• 강한 mortar → block sliding 어려움\n"
         "  → shear stiff (G +30%)\n"
         "• Hydrostatic은 vacancy void로 soft\n"
         "  → B −8%, but G ↑ dominant\n"
         "• E_VRH = 9BG/(3B+G) → +25%",
         size=11)

add_text(s, Inches(0.6), Inches(6.0), Inches(12.2), Inches(1.0),
         "※ B_VRH 23.4 ≈ B₀(BM-EOS) 21.7 — 두 독립 방법 cross-check (±3%)\n"
         "※ Kim ACS Mater. Lett. 2025 UPE 측정과 동일 방향 (Cl ↑ → E ↑) ✓",
         size=10, color=GRAY)
add_footer(s, 11)


# =====================================================
# SLIDE 8a — C44 shear lock-in
# =====================================================
s = slide_blank()
add_title(s, "8a. 어디서 +30% shear? — C44 +72% (4d-Cl shear lock-in)")

cij_data = [
    ["Cij (GPa, relaxed-ion)", "LPSCl", "LPSCl₁.₆", "Δ%"],
    ["C11", "37.7", "37.0", "−2%"],
    ["C12", "20.4", "16.8", "−18%"],
    ["C44 (shear) ★", "8.0", "13.7", "+72% !"],
    ["Zener A", "1.14", "1.44", "+26%"],
    ["ν (Poisson)", "0.36", "0.30", "−17%"],
]
add_table(s, Inches(0.6), Inches(1.4), Inches(6.5), Inches(2.9),
          cij_data, font_size=11, row_colors={3: LIGHT_RED})

add_rect(s, Inches(7.3), Inches(1.4), Inches(5.5), Inches(2.9),
         fill=LIGHT_GRAY)
add_text(s, Inches(7.4), Inches(1.5), Inches(5.3), Inches(0.4),
         "[ 4d-Cl shear pin cartoon ]",
         size=12, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(7.4), Inches(1.95), Inches(5.3), Inches(2.3),
         "comp1: Li 평면 슬라이드 자유\n"
         "modelc: 4d-Cl 짧고 강한 결합\n"
         "       (2.45 Å, ICOHP −2.84)이\n"
         "       특정 shear 방향 pin\n\n"
         "Hydrostatic은 모든 결합 균등\n"
         "압축 → pin 효과 안 켜짐 → B soft",
         size=11)

add_text(s, Inches(0.6), Inches(4.6), Inches(12.2), Inches(1.4),
         "• C44 +72% — VRH-평균 G +30%의 거의 모든 출처\n"
         "• C12 −18% — 추가 shear 자유도 감소\n"
         "• 4d-Cl AS (M3 ICOHP −2.84) 가 shear 방향 변형 선택적 lock\n"
         "• 등방 압축(C11 ~불변)엔 거의 영향 없음 → B₀ soft (vacancy void)",
         size=12)
add_footer(s, 12)


# =====================================================
# SLIDE 8b — Cross-check 4
# =====================================================
s = slide_blank()
add_title(s, "8b. M4 4 cross-check — paper-grade 정확도 보증")

cc_data = [
    ["#", "Cross-check", "LPSCl", "LPSCl₁.₆", "결론"],
    ["1", "B_VRH (stress-strain)\nvs B₀ (BM-EOS)", "25.5\n26.23", "23.4\n21.71", "±10% 일치 ✓"],
    ["2", "E_VRH (0K relaxed)\nvs Kim 2025 UPE / 자체 AFM", "22.06\n~12 GPa", "27.66\n14.9 GPa", "방향 일치 ✓"],
    ["3", "E_VRH (600 K MLIP snapshot)", "29.1", "32.9 (+13%)", "0K와 동일 방향 ✓"],
    ["4", "LOBSTER charge spilling\nk×L (paper-grade ≥40 Å)", "1.46%\n40", "1.16%\n42", "<5% / 수렴 ✓"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(4.5),
          cc_data, font_size=11, col_widths=[0.6, 4.2, 2.5, 2.5, 2.6])

add_text(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(0.7),
         "→ 4 독립 cross-check 모두 정합 — relaxed-ion +25% 결론 robust",
         size=14, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 13)


# =====================================================
# SLIDE 9 — Cross-check PS4 universal
# =====================================================
s = slide_blank()
add_title(s, "9. PS₄ covalent backbone — 두 시스템에서 사실상 동일")

ps4_data = [
    ["Probe", "LPSCl", "LPSCl₁.₆", "Δ"],
    ["P–S 길이 (Å, mean)", "2.073", "2.064", "−0.5%"],
    ["P–S σ (Å, 분산)", "0.036", "0.011", "modelc ↓ 더 균질"],
    ["P 배위수", "4.00", "4.00", "0% (완벽 보존)"],
    ["ICOHP P–S (eV/bond)", "−5.94", "−6.00", "+1.0%"],
    ["ELF P–S bridge", "0.946", "0.944", "Δ 0.002 (~0)"],
    ["Bader P (e)", "+4.69", "+4.43", "basin shape ¹"],
    ["Li–S(4d) universal anchor", "−2.57", "−2.52", "−2% (universal)"],
]
add_table(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.0),
          ps4_data, font_size=11,
          row_colors={7: LIGHT_BLUE})

add_text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.2),
         "• PS₄ 결합 길이 · ICOHP · ELF · 배위수 — 5 probe 모두 차이 +1% 이내\n"
         "• free S²⁻ 주위 Li 결합도 universal anchor (Δ2%)\n"
         "• 모든 조성 변화는 ionic sublattice 안에서만",
         size=12)

add_text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.3),
         "→ argyrodite의 PS₄³⁻ 단위 = chemistry-independent rigid block",
         size=13, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 14)


# =====================================================
# SLIDE 10 — Bond length per-site
# =====================================================
s = slide_blank()
add_title(s, "10. Li–Cl 평균 짧아짐 — 출처는 4d anti-site (per-bond vs per-anion)")

bond_data = [
    ["결합 (Å)", "LPSCl", "LPSCl₁.₆", "Δ"],
    ["P–S", "2.073 ± 0.036", "2.064 ± 0.011", "−0.5% (불변)"],
    ["Li–S", "2.461 ± 0.106", "2.465 ± 0.094", "+0.2% (동일)"],
    ["★ Li–Cl", "2.607 ± 0.129", "2.532 ± 0.119", "−3% 짧아짐 ↓"],
    ["S–S (cage)", "3.595", "3.519", "−2%"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(7.0), Inches(2.2),
          bond_data, font_size=10, row_colors={3: LIGHT_RED})

# Per-bond vs per-anion box
add_rect(s, Inches(7.7), Inches(1.3), Inches(5.2), Inches(2.2),
         fill=LIGHT_GRAY)
add_text(s, Inches(7.8), Inches(1.35), Inches(5.0), Inches(0.4),
         "Per-site (Li–Cl)", size=12, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(7.8), Inches(1.75), Inches(5.0), Inches(1.7),
         "Li–Cl (4a 정상): 2.59 Å (−1%)\n"
         "Li–Cl (4d AS):   2.45 Å (★ 0.14 Å 짧음)\n\n"
         "Voronoi V(Cl):\n"
         "  22.06 → 20.31 Å³ (−1.7)",
         size=11)

# Per-bond vs per-anion
pba_data = [
    ["측정 단위", "Li–S (4d free)", "Li–Cl (4d AS)", "결과"],
    ["per-bond ICOHP (eV)", "−2.57", "−2.84", "Cl 약간 ↑ (LOBSTER overlap)"],
    ["per-anion total (eV)\n= per-bond × coord", "−2.57 × 6 = −15.4", "−2.84 × 4 = −11.3", "S²⁻ 36% ↑ (Coulomb 회복)"],
]
add_table(s, Inches(0.5), Inches(3.7), Inches(12.4), Inches(2.0),
          pba_data, font_size=10)

add_text(s, Inches(0.5), Inches(5.9), Inches(12.4), Inches(1.0),
         "• Per-bond: LOBSTER가 covalent overlap 우세 측정 → Cl 약간 우세\n"
         "• Per-anion total: S²⁻ q² × coord(6 vs 4) → Coulomb 직관 회복\n"
         "• 4d cage 자체가 짧은 Li 거리 (격자 기하) — ionic radius 단독으로 못 설명",
         size=11)
add_footer(s, 15)


# =====================================================
# SLIDE 11 — Voronoi 4-sublattice
# =====================================================
s = slide_blank()
add_title(s, "11. Voronoi 4-Sublattice — Disorder가 어느 sublattice에 들어가나")

vor_data = [
    ["원자", "comp1 std (Å³)", "modelc std (Å³)", "변화", "역할"],
    ["P (PS₄ framework)", "0.00", "0.37", "+0.37", "framework 거의 불변 (M1)"],
    ["Cl", "0.00", "0.74", "+0.74", "4a + 4d 두 군집 분기"],
    ["★ Li", "0.21", "1.15", "×5.5 ↑", "disorder fingerprint #1"],
    ["★ S", "3.41", "2.05", "−40% ↓", "역설적 균질화"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(2.6),
          vor_data, font_size=11,
          row_colors={3: LIGHT_RED, 4: LIGHT_RED})

# Interpretation panel
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.4), Inches(2.0),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.3), Inches(12.2), Inches(0.4),
         "4 sublattice 응답 해석",
         size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.7), Inches(4.7), Inches(12.2), Inches(1.4),
         "• P  → PS₄ framework 거의 안 흔들림 (M1 backbone 일관)\n"
         "• Cl → 4a (정상) + 4d AS 두 환경 분기 (M3 ICOHP per-site)\n"
         "• Li → ×5.5 흔들림 — vacancy + anti-site 인접 재배치\n"
         "• S  → 오히려 균질화 — anti-site Cl이 PS₄-S vs free-S²⁻ split 메움",
         size=12)

add_text(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.5),
         "→ 3 독립 probe 수렴 (Voronoi + BVSE bimodal + LOBSTER ICOHP +13%)",
         size=13, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 16)


# =====================================================
# SLIDE 12 — BVSE bimodal
# =====================================================
s = slide_blank()
add_title(s, "12. BVSE Bimodal Split — Li 환경 둘로 갈라짐 (paired 5×5×5)")

bvse_data = [
    ["시스템", "Group", "n_Li", "비율", "BVS peak", "환경"],
    ["comp1", "uniform", "3000", "100%", "1.60–1.64", "F-43m ordered"],
    ["modelc", "low-BVS (A)", "1074", "39.8%", "1.60–1.64", "comp1-like, AS 멀리"],
    ["modelc ★", "high-BVS (B)", "1626", "60.2%", "1.83–1.89", "AS Cl 인접, +15%"],
]
add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(2.2),
          bvse_data, font_size=11, row_colors={3: LIGHT_ORANGE})

# Cross-check box
add_rect(s, Inches(0.4), Inches(3.8), Inches(12.6), Inches(2.2),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.6), Inches(3.9), Inches(12.2), Inches(0.4),
         "정량 cross-validation",
         size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.6), Inches(4.3), Inches(12.2), Inches(1.6),
         "• 1626 high-BVS Li / 300 AS Cl  =  5.4 Li per AS (Li–Cl coord 4-6 정합)\n"
         "• 60.2% 비율 = 5.4 × 300 / 2700 (math 정확)\n"
         "• BVSE +15% shift  ↔  LOBSTER ICOHP +40% per-bond (slide 7)\n"
         "    같은 anti-site 효과를 두 독립 probe가 다른 정량으로 측정",
         size=12)

add_text(s, Inches(0.4), Inches(6.2), Inches(12.6), Inches(0.7),
         "※ 37.5% AS = cubic 5×5×5 stoichiometric necessity (Cl ≥ 1.6/fu, 4a 500개 full)\n"
         "  실험 25–50% 범위 내 — 'over-disordered' 아님",
         size=10, color=GRAY)
add_footer(s, 17)


# =====================================================
# SLIDE 13 — ELF
# =====================================================
s = slide_blank()
add_title(s, "13. ELF — Covalent backbone vs Ionic glue 시각화")

# Two ELF placeholders
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.5), Inches(2.6), Inches(6.0), Inches(0.6),
         "[ ELF 2D slice — comp1 ]\nPS₄ 적색 maxima (0.95)\nLi 청색 depletion (<0.1)",
         size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(6.7), Inches(1.4), Inches(6.0), Inches(3.5),
         fill=LIGHT_GRAY)
add_text(s, Inches(6.7), Inches(2.6), Inches(6.0), Inches(0.6),
         "[ ELF 2D slice — modelc ]\n동일 패턴 (PS₄ 0.94)\n4d-Cl AS 추가",
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

elf_data = [
    ["Probe location", "LPSCl", "LPSCl₁.₆", "의미"],
    ["P–S bond midpoint", "0.946", "0.944", "covalent 동일 ★"],
    ["Li basin floor", "0.072", "0.065", "ionic depletion 강함"],
    ["Li → nearest anion line min", "0.07", "0.04", "modelc 더 ionic"],
]
add_table(s, Inches(0.5), Inches(5.05), Inches(12.4), Inches(1.55),
          elf_data, font_size=11, row_colors={1: LIGHT_BLUE})

add_text(s, Inches(0.5), Inches(6.68), Inches(12.4), Inches(0.3),
         "→ PDOS + LOBSTER + ELF 세 독립 probe — 같은 그림 (covalent + ionic 공존)",
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 18)


# =====================================================
# SLIDE 14 — k-mesh audit
# =====================================================
s = slide_blank()
add_title(s, "14. k-mesh Audit + LOBSTER spilling + AIMD window")

# k-mesh incident box
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(2.0),
         fill=LIGHT_RED, line=HIGHLIGHT_RED)
add_text(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.4),
         "comp1 k-mesh incident & recovery — 정직한 referee defense",
         size=14, bold=True, color=HIGHLIGHT_RED)
add_text(s, Inches(0.7), Inches(1.95), Inches(12.0), Inches(1.4),
         "초기:  k = 2×2×1, k×L = 10 Å  →  gap 1.50 eV (artifact)\n"
         "발견:  Δgap (modelc 1.82) = 0.32 eV — '조성 효과'로 잘못 해석\n"
         "복구:  k = 4×4×4, k×L = 40 Å  →  gap 1.76 eV (paper-grade)\n"
         "결과:  Δgap = 0.06 eV (M1 확정), structure RMS 0.003 Å (electronic만 오염)",
         size=12)

ksens = [
    ["Property", "k-sens", "대응"],
    ["B₀ (BM-EOS)", "robust", "재계산 불필요 (volume curvature systematic cancel)"],
    ["Gap", "high", "k = 4×4×4 / 6×6×3 재계산 ✓"],
    ["Elastic Cij", "high", "재계산 ✓"],
    ["DOS shape", "high", "재계산 ✓"],
    ["ICOHP", "robust", "Δ < 0.006 eV (local probe) ✓"],
]
add_table(s, Inches(0.5), Inches(3.7), Inches(12.4), Inches(2.3),
          ksens, font_size=10)

add_text(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(0.7),
         "LOBSTER spilling: comp1 1.46% / modelc 1.16% (paper <5% ✓)\n"
         "AIMD window: [2, 50] ps 양쪽 동일 protocol, R² 0.999/0.992",
         size=11, color=GRAY)
add_footer(s, 19)


# =====================================================
# SLIDE 16 — Summary (15 dropped)
# =====================================================
s = slide_blank()
add_title(s, "16. Summary — 불변 ↔ 변화의 깔끔한 분리")

# Two columns
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(4.7),
         fill=LIGHT_BLUE, line=HANYANG_BLUE, line_w=1)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.4),
         "불변 (조성 둔감)", size=16, bold=True, color=HANYANG_BLUE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.65), Inches(1.85), Inches(5.7), Inches(4.0),
         "• PS₄ covalent backbone\n"
         "  P–S length, ICOHP, ELF\n\n"
         "• Band gap (Δ 0.06 eV)\n"
         "  VBM/CBM character\n\n"
         "• Li–S(4d) universal anchor\n"
         "  (Δ 2%)\n\n"
         "• Per-anion Coulomb 직관\n"
         "  (S²⁻ 우세)",
         size=12)

add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(4.7),
         fill=LIGHT_RED, line=HIGHLIGHT_RED, line_w=1)
add_text(s, Inches(6.95), Inches(1.4), Inches(5.8), Inches(0.4),
         "변화 (조성 민감)", size=16, bold=True, color=HIGHLIGHT_RED,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.0),
         "• Li transport: Ea ↓ + D₀ ↑\n"
         "  σ × 2.5–4.3 (T 의존)\n\n"
         "• Li–anion ionic glue\n"
         "  ICOHP +13% (Cl), +8% (S)\n\n"
         "• Shear modulus G (+30%)\n"
         "  C44 +72%, Zener A 1.14→1.44\n\n"
         "• 4d-Cl anti-site (새 결합 family)\n"
         "  ICOHP −2.84 eV (+40%/bond)",
         size=12)

# Thesis box
add_rect(s, Inches(0.5), Inches(6.15), Inches(12.4), Inches(0.85),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.4), Inches(0.4),
         '"전자구조가 아니라 구조적 무질서 (Li 공공 + 4d-Cl anti-site)에서 온다"',
         size=14, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.4), Inches(0.4),
         "→ covalent skeleton 유지 + ionic ligament 재배치",
         size=12, color=BLACK, align=PP_ALIGN.CENTER, font="맑은 고딕")
add_footer(s, 20)


# =====================================================
# SLIDE 17 — 3-probe convergence
# =====================================================
s = slide_blank()
add_title(s, "17. 3-Probe Convergence — Anti-site Cl 효과 같은 그림")

probe_data = [
    ["Probe", "정량 결과", "측정 대상", "Slide"],
    ["Voronoi V std", "Cl 0 → 0.74 Å³\nLi 0.21 → 1.15 (×5.5)", "Site disorder (geometric)", "11"],
    ["BVSE bimodal", "+15% shift (60.2% Li)\n1.62 → 1.85 BVS", "Li 이동 환경 (path)", "12"],
    ["LOBSTER ICOHP", "+40% per-bond (4d AS)\n−2.03 → −2.84 eV", "Li-Cl 결합 강도 (chemistry)", "7"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(3.0),
          probe_data, font_size=11)

# Cross-consistency box
add_rect(s, Inches(0.5), Inches(4.6), Inches(12.4), Inches(2.0),
         fill=LIGHT_ORANGE)
add_text(s, Inches(0.7), Inches(4.7), Inches(12.0), Inches(0.4),
         "Quantitative cross-consistency",
         size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(1.4),
         "• 1626 high-BVS Li / 300 AS Cl = 5.4 Li/AS  ↔  Li–Cl coord (4-6) 정합\n"
         "• BVSE 60.2% group = 5.4 × 300 / 2700 (math 정확)\n"
         "• Δ(BVS) +15% vs Δ(ICOHP) +40% — 다른 measure same anchor",
         size=12)

add_text(s, Inches(0.5), Inches(6.8), Inches(12.4), Inches(0.4),
         "→ 단일 method artifact 아닌 robust 물리 현상",
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 21)


# =====================================================
# SLIDE 18 — Oxidation 4-axis
# =====================================================
s = slide_blank()
add_title(s, "18. Oxidation — 'Cl-rich 더 안정?' 축을 지정해야 답")

ox_data = [
    ["Axis", "의미", "comp1 vs modelc", "Reference"],
    ["1", "0-pressure intrinsic\nbulk redox onset", "DRAW ~2.14 V\n(S²⁻ limited)", "우리 + Gil 2022"],
    ["2 ★", "Mechanically constrained\n(K_eff > 0)", "Cl-rich WINS\nwindow 0.80–4.30 V", "Gil-González ESM 2022"],
    ["3 ★", "Cathode interface cycling\n(R_int, CE)", "Cl-rich WINS\n8.9 vs 13.2 Ω·h^-0.5", "Zuo Angew 2023"],
    ["4 ✗", "Thermal / calendar aging\n(shelf life 90°C)", "Cl-rich LOSES\n48% vs 68% (5d)", "Wu Nano En 2026"],
]
row_colors_ox = {2: LIGHT_GREEN, 3: LIGHT_GREEN, 4: LIGHT_RED}
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(4.5),
          ox_data, font_size=11, row_colors=row_colors_ox)

add_text(s, Inches(0.5), Inches(6.05), Inches(12.4), Inches(0.7),
         "→ modelc σ 향상이 oxidation penalty 없이 옴 (axis 1–3 neutral or favorable)\n"
         "   비용은 thermal shelf life (axis 4) — paper #2 motivation",
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 22)


# =====================================================
# SLIDE 19 — ESW + decomp
# =====================================================
s = slide_blank()
add_title(s, "19. Constrained ESW Cl-scan + 분해반응 — Axis 1–3 정량 backing")

esw_data = [
    ["K_eff (GPa)", "Cl 0.5", "Cl 1.0", "Cl 1.5", "Cl 1.6 (modelc)", "Cl 2.0"],
    ["0", "2.40", "2.40", "2.40", "2.40 (flat)", "2.40"],
    ["10", "3.20", "3.40", "3.80", "4.05 ★", "3.50 ↓"],
    ["20", "3.70", "4.10", "4.20", "4.30 ★★", "3.60 ↓↓"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(2.0),
          esw_data, font_size=11, row_colors={3: LIGHT_GREEN})

# Decomposition box
add_rect(s, Inches(0.5), Inches(3.6), Inches(12.4), Inches(2.2),
         fill=LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(3.65), Inches(12.0), Inches(0.4),
         "분해반응 비교 (oxidation onset 2.14 V)",
         size=13, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.7), Inches(4.05), Inches(12.0), Inches(1.7),
         "comp1:   Li₆PS₅Cl  →  Li₃PS₄ + 0.25 LiS₄ + LiCl + 1.75 Li ↑\n"
         "modelc:  Li₅.₄PS₄.₄Cl₁.₆  →  Li₃PS₄ + 0.1 LiS₄ + 1.6 LiCl + 0.7 Li ↑\n\n"
         "modelc: +1.6× LiCl (더 inert) · −2.3× Li 방출\n"
         "→ Zuo 2023 Eq(1)/(2) 정량 cross-validation ✓",
         size=12)

add_text(s, Inches(0.5), Inches(6.0), Inches(12.4), Inches(0.7),
         "→ Cl-rich가 axis 1–3에서 우세하다는 결과를 우리 직접 계산이 정량 backing",
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 23)


# =====================================================
# SLIDE 20 — Tension audit
# =====================================================
s = slide_blank()
add_title(s, "20. 4가지 Literature Tension — 모두 해소")

ten_data = [
    ["#", "Tension", "우리 결과", "해소"],
    ["1", "Oxidation: Cl-rich 더 안정?\n(lit consensus: yes)",
     "0-pressure DRAW\n(slide 18 axis 1)",
     "4-axis framework\n(slide 18–19)"],
    ["2", "Ea-vs-prefactor mechanism\n(Minafra: disorder→Ea↓)",
     "comp1 4fu Ea > modelc\nMinafra direction 정합 ★",
     "Schlem 정확 매칭\n(slide 6 v3)"],
    ["3", "Elastic vacancy paradox\n(Kim 2025: Cl → E up)",
     "clamped: 52 vs 52 (동일)\nrelaxed: +25% (Kim 매칭)",
     "Relaxed-ion 정확한\nmethod (slide 8)"],
    ["4", "Band gap absolute\n(PBE PAW lit ~2.3)",
     "PBE/USPP 1.76 vs lit 2.3",
     "Method offset 양쪽\n동일, Δgap robust"],
]
add_table(s, Inches(0.3), Inches(1.4), Inches(12.7), Inches(4.7),
          ten_data, font_size=10, col_widths=[0.6, 4.4, 4.0, 3.7])

add_text(s, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.6),
         "→ 4 tension 모두 'artifact 아닌 method/축 분리'로 해소\n"
         "   paper의 nuanced honesty 핵심 슬라이드",
         size=12, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 24)


# =====================================================
# SLIDE 21 — Caveats
# =====================================================
s = slide_blank()
add_title(s, "21. 모든 Caveat 한 페이지 — 정직한 method limitation")

cav_data = [
    ["#", "Caveat", "우리 대응"],
    ["1", "UMA MLIP σ overshoot ~3–5×", "σ ratio만 사용 (절대 σ 인용 X)"],
    ["2", "Haven ratio = 1 가정", "σ_NE는 upper bound로 표기"],
    ["3", "3-pt Arrhenius 300K 외삽", "정성적 trend만, 정량 단정 X"],
    ["4", "k-mesh convergence 필요", "k×L ≥ 40 Å 양쪽 보장 (slide 14)"],
    ["5", "LOBSTER charge spilling", "1.16% / 1.46% (<5% paper 기준)"],
    ["6", "5 f.u. cell / n=3 configs", "ΔEa < 0.05 eV 해상도 못 잡음"],
    ["7", "USPP gap 0.4 eV underestimate", "method offset 명시, Δgap robust"],
    ["8", "0-pressure ESW only", "K_eff axis 분리 (slide 18–19)"],
    ["9", "Random anti-site disorder model", "실험 charge-coupled placement 다를 가능성"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(4.7),
          cav_data, font_size=10, col_widths=[0.6, 5.6, 6.2])

add_text(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.6),
         "→ 어떤 caveat도 4 메시지 자체를 흔들지 못함\n"
         "   trend / ratio / mechanism은 robust (same-protocol paired comparison의 힘)",
         size=12, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 25)


# =====================================================
# SLIDE 22 — Trade-offs + Outlook
# =====================================================
s = slide_blank()
add_title(s, "22. Trade-offs & Outlook — Paper #2 Doping Strategy로의 다리")

# Top panel: Trade-offs
add_rect(s, Inches(0.5), Inches(1.3), Inches(12.4), Inches(2.4),
         fill=LIGHT_RED, line=HIGHLIGHT_RED, line_w=1)
add_text(s, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.4),
         "4 Trade-offs of LPSCl₁.₆ (idle/storage 영역만)",
         size=14, bold=True, color=HIGHLIGHT_RED)
add_text(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(1.9),
         "1. Thermal calendar shelf life — Wu 2026: 90°C 5일, L6 68% > L55 48% > L53(Cl1.7) 35%\n"
         "2. Moisture sensitivity — Cl⁻ + H₂O → LiOH·LiCl + H₂S (Strauss, Kraft)\n"
         "3. Synthesis window — Cl ≥ 1.7 phase-pure 어려움 (Adeli, Yu, Wu)\n"
         "4. Mechanical anisotropy — Zener A 1.14 → 1.44 (mild but cycling fatigue 가능)\n\n"
         "Common cause: 4d-Cl anti-site disorder (M3-M4의 source가 동시에 trade-off의 source)",
         size=11)

# Arrow
add_text(s, Inches(0.5), Inches(3.85), Inches(12.4), Inches(0.3),
         "↓  Paper #2 Strategy  ↓",
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)

# Bottom panel: Mitigation
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.4), Inches(2.4),
         fill=LIGHT_BLUE, line=HANYANG_BLUE, line_w=1)
add_text(s, Inches(0.7), Inches(4.25), Inches(12.0), Inches(0.4),
         "Mitigation via Oxide Doping (X₂O_y → SE)",
         size=14, bold=True, color=HANYANG_BLUE)
add_text(s, Inches(0.7), Inches(4.65), Inches(12.0), Inches(1.9),
         "Hypothesis: O²⁻ PS₄ 부분 치환 또는 4d 자리 점유 → O–Li 강한 결합, Cl/O mixed sublattice, LiCl 2차상 차단\n\n"
         "Active 후보 (14 dopant × 3 conc cascade screening):\n"
         "  ★ Sc₂O₃   — cascade strongest (de_post = −0.974, E_VRH 18.7 GPa)\n"
         "  ★ B₂O₃    — thermal stabilizer, anneal+EOS 진행\n"
         "  ★ Nd₂O₃   — DFT-relaxed run5 완료, EOS+post 대기\n"
         "  • Al₂O₃ cluster · MnO/CoO 부드러움 후보 (41 champions verified)",
         size=11)

# Closing
add_text(s, Inches(0.5), Inches(6.75), Inches(12.4), Inches(0.5),
         '"구조적 무질서의 양면성을 이해했으니, 이제 그 무질서를 \'design\'할 수 있다"',
         size=12, bold=True, color=HANYANG_BLUE, align=PP_ALIGN.CENTER)
add_footer(s, 26)


# =====================================================
# Figure-attachment markers (사진 첨부 표시)
# =====================================================
FIG_YELLOW = RGBColor(0xFF, 0xF1, 0x9A)
FIG_BORDER = RGBColor(0xC8, 0x9A, 0x00)
FIG_TEXT = RGBColor(0x66, 0x4D, 0x00)


def add_fig_chip(slide, label, where="title", w=5.6):
    """Yellow chip indicating a figure should be inserted.
    where='title' → top-right corner, ABOVE the title text so long titles
    are never covered. where='bottom' → just above footer."""
    if where == "title":
        left = Inches(13.333 - w - 0.15)
        top = Inches(0.03)
    else:
        left = Inches(13.333 - w - 0.5)
        top = Inches(6.55)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, Inches(w), Inches(0.30))
    shp.fill.solid()
    shp.fill.fore_color.rgb = FIG_YELLOW
    shp.line.color.rgb = FIG_BORDER
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"📎 사진 첨부: {label}"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = FIG_TEXT
    r.font.name = "맑은 고딕"


def add_fig_box(slide, left, top, width, height, label):
    """Large dashed placeholder box with caption, for slides where a figure
    is the main content. Drawn dashed-yellow."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xE6)
    shp.line.color.rgb = FIG_BORDER
    shp.line.width = Pt(1.5)
    ln = shp.line._get_or_add_ln()
    pr = ln.makeelement(qn("a:prstDash"), {"val": "dash"}, nsmap=None)
    ln.append(pr)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"📎 [ 사진/그림 자리 ]\n{label}"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = FIG_TEXT
    r.font.name = "맑은 고딕"


# Slide order (0-based) → (label, position)
# Title slide (0) and pure-text summary slides skipped.
FIG_MARKERS = {
    1:  ("결정구조 VESTA (F-43m vs R3m)",        "title"),
    2:  ("3-tier pipeline 모식도 (선택)",         "title"),
    3:  ("4-message at-a-glance figure",          "title"),
    4:  ("PDOS overlay (comp1 vs modelc)",        "title"),
    5:  ("σ(T) Arrhenius plot (3점 paired)",      "title"),
    6:  ("D₀ vs Eₐ 산점도 / prefactor 막대",      "title"),
    7:  ("Disorder ensemble Eₐ 분포",             "title"),
    8:  ("σ(T) 저온 외삽 plot (300 K)",           "title"),
    9:  ("Li–anion ionic glue 막대/모식도",        "title"),
    10: ("B·G·E clamped vs relaxed 막대",         "title"),
    11: ("Cij component bar (C44 강조)",          "title"),
    12: ("M4 cross-check 수렴 그림",              "title"),
    13: ("PS₄ 결합 히스토그램 overlay",            "title"),
    14: ("Li–Cl per-bond 히스토그램 (4a/4d)",     "title"),
    15: ("Voronoi 4-sublattice 모식도",           "title"),
    16: ("BVSE bimodal 분포 / pathway 맵",        "title"),
    17: ("ELF isosurface (covalent+ionic)",       "title"),
    18: ("k-mesh / spilling / AIMD 수렴",         "title"),
    20: ("3-probe 수렴 schematic",                "title"),
    21: ("Oxidation 4-axis radar/grid",           "title"),
    22: ("ESW Cl-scan + 분해반응 ΔG 막대",        "title"),
    25: ("Trade-off 4축 + 도핑 roadmap",          "title"),
}

slide_list = list(prs.slides)
for idx, (label, where) in FIG_MARKERS.items():
    if idx < len(slide_list):
        add_fig_chip(slide_list[idx], label, where=where)


# =====================================================
# Save
# =====================================================
out = "/home/user/Yonghoon-DEM-DFT/kb/papers/lpscl_vs_lpscl16_seminar_v1.pptx"
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides)}")
print(f"fig markers: {len(FIG_MARKERS)}")
