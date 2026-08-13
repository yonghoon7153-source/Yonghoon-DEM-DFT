#!/usr/bin/env python3
"""revise_cascade_deck.py — Research Seminar 2026-08 cascade 덱 개정.

기존 릴리스 덱(28장)에 본문 2장을 추가하고 4장을 고친다. 새 장은 **기존 장을 deepcopy 해서
텍스트만 갈아끼우는** 방식이라 폰트·색·여백이 원본과 100% 같다 (좌표를 손으로 재현하지 않는다).

추가:
  P12b  외부 대조 — 방법 원저자 그룹(Banik/Zeier/Mo)이 같은 질문을 했다
  P12c  트레이드오프 2 — 값 하나가 두 축을 갈라놓았다 (E_VRH 순위 재계산)
수정:
  P18   철회 원장에 5번째 줄 (교차파일 값 충돌)
  P21   마무리를 '트레이드오프 2건 관측'으로 승격
  A6/A7 출처 행 추가

이 도구가 못 하는 것: 그림을 새로 그리지 않는다(텍스트 슬라이드만 만든다).
레이아웃 자동 조정도 안 한다 — 행이 늘면 pitch 를 코드에서 직접 지정해야 한다.
"""
import argparse, copy, sys
from pptx import Presentation
from pptx.util import Inches, Pt


def clone_after(prs, src_idx, dst_idx):
    """src_idx(1-base) 슬라이드를 복제해 dst_idx(1-base) 위치에 넣는다."""
    src = prs.slides[src_idx - 1]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    xml = prs.slides._sldIdLst
    ids = list(xml)
    xml.remove(ids[-1])
    xml.insert(dst_idx - 1, ids[-1])
    return prs.slides[dst_idx - 1]


def texts(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]


def set_text(shape, s):
    """첫 run 의 서식을 유지한 채 문자열만 교체."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    keep = p0.runs[0]
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    keep.text = s
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def find(slide, needle):
    for sh in texts(slide):
        if needle in sh.text_frame.text:
            return sh
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args()
    if a.selftest:
        sys.exit(selftest())

    prs = Presentation(a.src)
    n0 = len(prs.slides)
    if n0 != 28:
        print(f"⚠ 예상 28장이 아니라 {n0}장 — 위치 지정이 어긋날 수 있다")

    # ---------- P12b : 외부 대조 (slide 18 격자를 복제) ----------
    s = clone_after(prs, 18, 13)
    set_text(find(s, "Validation"), "External check")
    set_text(find(s, "Retraction ledger"), "Same method, same question, independent group")
    set_text(find(s, "Failures made the cascade more credible"),
             "The originators of our oxidation method asked whether substitution can move the onset")
    rows = ["Single-seed 1.33", "MSD fit outside", "DOS-threshold", "raw HSAB grade"]
    new_rows = [
        "Banik / Zeier / Mo (2022) — same grand-potential tool + HAXPES + stepwise CV",
        "Their answer:  VBM is pinned by S (free S²⁻ and non-bonding PS₄³⁻ 3p) — substitution does not move it",
        "Our screen:  six candidates raise the onset — and all six then stop at the Li-path gate",
        "Same conclusion in different words:  the oxidation axis is not bought cheaply",
    ]
    for old, new in zip(rows, new_rows):
        set_text(find(s, old), new)
    set_text(find(s, "Failed claims became"),
             "Their substitutions stay inside the sulfide framework; ours add oxide / fluoride units — different class, same verdict shape.")
    set_text(find(s, "Source:"),
             "Source: litdb banik2022_substitutions_oxidative_stability_argyrodite · Zeier (Münster) + Mo (Maryland)")

    # ---------- P12c : 트레이드오프 2 ----------
    s = clone_after(prs, 19, 14)
    set_text(find(s, "Validation"), "Result")
    set_text(find(s, "Retraction ledger"), "Trade-off 2 · found while verifying a value")
    set_text(find(s, "Failures made the cascade more credible"),
             "One corrected number split two axes that used to agree")
    new_rows = [
        "Old record:  Sc₂O₃ x002 — softest (E_VRH 18.7 GPa) AND strongest formation (−0.974 eV/atom)",
        "Cross-file conflict traced to raw output:  E_VRH is 42.082, not 18.7 — a factor of 2.2",
        "Re-sorted all 40 champions:  softest are now Li₂O 32.4 · MnO 32.6 · NiO 33.9;  Sc₂O₃ falls to rank 18",
        "Formation energy stays rank 1 (−0.9338) — so the two axes came apart",
    ]
    for old, new in zip(rows, new_rows):
        set_text(find(s, old), new)
    set_text(find(s, "Failed claims became"),
             "Gate audit showed one trade-off; value audit showed another. In this dataset, one candidate winning every axis does not survive checking.")
    set_text(find(s, "Source:"),
             "Source: doping_cascade_verified.json (recomputed 2026-08-13) · gabia FINAL_RANKING.json raw rows")

    # ---------- P18 : 철회 원장 5번째 줄 ----------
    s = prs.slides[19]                      # 두 장 삽입 후 18 → 20번째(0-base 19)
    ledger = [find(s, k) for k in rows]
    boxes = [sh for sh in s.shapes if not sh.has_text_frame or not sh.text_frame.text.strip()]
    tops = [Inches(1.62 + 0.84 * i) for i in range(5)]
    for i, sh in enumerate(ledger):
        sh.top = tops[i] + Inches(0.18)
        sh.height = Inches(0.46)
    fifth = copy.deepcopy(ledger[-1]._element)
    s.shapes._spTree.append(fifth)
    fifth_sh = s.shapes[-1]
    fifth_sh.top = tops[4] + Inches(0.18)
    set_text(fifth_sh, "Cross-file value conflict   →   raw-output arbitration   →   ranking recomputed (claim changed, not a digit)")
    # 배경 박스도 같은 pitch 로 재배치 + 한 장 추가
    bg = [sh for sh in s.shapes
          if not (sh.has_text_frame and sh.text_frame.text.strip())
          and Inches(0.7) < sh.left < Inches(0.9) and sh.height > Inches(0.5)]
    for i, sh in enumerate(bg[:4]):
        sh.top = tops[i]; sh.height = Inches(0.78)
    if bg:
        nb = copy.deepcopy(bg[0]._element)
        s.shapes._spTree.insert(0, nb)
        s.shapes[0].top = tops[4]; s.shapes[0].height = Inches(0.78)

    # ---------- P21 마무리 ----------
    # 부제(Sendek/Kim 패턴)는 건드리지 않는다 — 제목줄과 중복되면 같은 말이 두 번 뜬다.
    s = prs.slides[22]
    tail = find(s, "physics and experiments decide")
    if tail:
        set_text(tail, "Two independent trade-offs say no single winner exists — "
                       "ML chooses what to calculate next; physics and experiments decide what is true.")

    # ---------- 페이지 번호 재부여 ----------
    for i, sl in enumerate(prs.slides, 1):
        for sh in texts(sl):
            t = sh.text_frame.text.strip()
            if t.isdigit() and sh.left > Inches(8.9):
                set_text(sh, str(i))

    prs.save(a.out)
    print(f"-> {a.out}  ({n0} → {len(prs.slides)}장)")


def selftest():
    ok = True
    def chk(n, c):
        nonlocal ok
        print(("  ok   " if c else "  FAIL ") + n); ok = ok and bool(c)
    from pptx import Presentation as P
    import io
    prs = P()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb.text_frame.text = "alpha"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(19)
    set_text(tb, "beta")
    chk("set_text 가 문자열을 바꾼다", tb.text_frame.text == "beta")
    chk("set_text 가 서식을 유지한다", tb.text_frame.paragraphs[0].runs[0].font.size == Pt(19))
    prs.slides.add_slide(prs.slide_layouts[6])
    n = len(prs.slides)
    clone_after(prs, 1, 2)
    chk("clone_after 가 장수를 늘린다", len(prs.slides) == n + 1)
    chk("clone_after 가 지정 위치에 넣는다", prs.slides[1].shapes[0].text_frame.text == "beta")
    chk("원본이 남아 있다", prs.slides[0].shapes[0].text_frame.text == "beta")
    chk("find 가 없는 문자열에 None", find(prs.slides[0], "zzz") is None)
    print("RESULT:", "0 실패" if ok else "실패 있음")
    return 0 if ok else 1


if __name__ == "__main__":
    main()
