#!/usr/bin/env python3
"""build_seminar_slide.py — 세미나 덱의 한 장을 단독 파일로 다시 만든다.

  python3 tools/seminar/build_seminar_slide.py 8
  python3 tools/seminar/build_seminar_slide.py 9
  python3 tools/seminar/build_seminar_slide.py --selftest

왜 한 도구인가 — 8·9장이 **같은 머리말 골격**(도형 0–13)을 쓴다. 장마다 파일을
  만들면 레이아웃 수정을 N 번 해야 하고, 실제로 8장에서 고친 세 가지(용어줄 위치·
  중복 각주·글자폭)가 9장에 그대로 남아 있었다. 문장만 SLIDES 에 둔다.

8장 배경 (2026-08-18)

왜 (2026-08-18)
  v6 덱의 8장에 세 가지 문제가 있었다. 전부 **화면에서 눈으로 보이는** 것들이다:
    ① 용어 설명줄(shape 2)이 헤더 불릿(shape 5) 위에 올라타고 오른쪽이 잘렸다
       — `how many compounds of that eleme` 에서 끊겼다.
    ② 같은 단서("superseded snapshot")가 **세 번** 나왔다: 소불릿 · 그림 밑 빨간 줄 ·
       맨 아래 이탤릭. 그 이탤릭도 잘려서 `current rank ing is 0 species` 로 보였다.
    ③ 색이 47종 풀(2026-06-29)이었다. 89종 풀로 바꾸면서 캡션도 같이 가야 한다.

  ⚠ 소불릿의 "superseded snapshot" 은 **변명**이었다. 89종으로 다시 매겨도 상·하위
    집합이 그대로라는 실측이 나왔으므로(카드 아래 참조), 이제 근거 있는 문장으로 바꾼다.

  python3 tools/seminar/build_slide08_periodic.py --selftest
  python3 tools/seminar/build_slide08_periodic.py

⚠ PowerPoint 자동서식 주의 (2026-08-18 실측)
  빌더가 쓴 것은 평범한 run 이라 그대로 나온다. 그런데 **PowerPoint 안에서 다시 타이핑**하면
  `A 2+ dopant` 의 `2+` 가 위첨자로 올라가 **A²⁺**(원소 A 의 2가 이온)처럼 보인다.
  → 이 파일에서 만든 슬라이드는 **복사해 붙이고, 그 줄을 다시 치지 않는다.**
     (1저자 판단 2026-08-18: 문구는 `A 2+` 로 유지하고 입력 쪽에서 처리한다.)

이 도구가 못 하는 것
  · 점수를 다시 매기지 않는다. 그림은 plot_seminar_2026_08.fig_periodic 이 만든다.
  · 원본 덱(v6)을 고치지 않는다 — 단독 파일만 만든다. 덱 반영은 사람이 붙여넣는다.
  · 머리말 도형(0–13)의 위치를 바꾸지 않는다. 로고·푸터 규약을 건드리면 안 된다.
"""
import os
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.insert(0, os.path.join(ROOT, "tools", "seminar"))

SRC = os.path.join(ROOT, "docs", "Research_Seminar_2026_08_cascade_story_v6.pptx")
SCRIPT = os.path.join(ROOT, "kb", "seminars", "cascade_dopant_screening_story_2026_08.md")

#: 장별 문장 — 대본(kb/seminars/…)과 **1:1**이어야 한다. selftest 가 대조한다.
#: idx 는 0-based (덱의 N 장 = idx N-1). fig 는 repo 루트 기준 (add_pics 규약).
SLIDES = {
 8: dict(
    idx=7, out="Slide08_periodic_table.pptx",
    fig="docs/figures/seminar/roster_periodic_table.png",
    zone=dict(x=0.55, y=2.28, w=8.90, h=3.62),
    title="Candidate set on the periodic table",
    head="Candidate set (2): 36 cation elements, colored by their best score",
    sub1="Late transition metals (Fe–Cu) form the [r]red[/r] block — their oxidation window is the narrowest.",
    sub2="[b]Group trends[/b] are robust, the exact order is not — the same elements stay at both ends.",
    label="best score per element, 89-species pool (2026-08-13)",
    note="[r]no approved ranking yet[/r] — the same elements stayed top and bottom when the set doubled (47 → 89)",
    gloss=("Cation: the metal we substitute in  ·  Coverage: how many compounds of that element "
           "were scored  ·  de: stability against the host"),
    drop=[10],
 ),
 9: dict(
    idx=8, out="Slide09_site_and_charge.pptx",
    fig="docs/figures/seminar/step1_site_choice.png",
    zone=dict(x=0.55, y=2.28, w=8.90, h=3.62),
    title="Substitution site and charge compensation",
    head="Step 1: The substitution site is enumerated before a structure is built",
    sub1="Li, P, S and Cl are separate sublattices — an oxide has to take [b]one cation and one anion site[/b].",
    sub2="A 2+ dopant on a Li site leaves extra charge — [r]each way of balancing it is a different structure[/r].",
    label="site preference vs ionic radius",
    note="only Si always takes the P framework; 19 of 26 always take Li — [r]6 switch with the doping level[/r]",
    gloss=("Sublattice: equivalent positions of one element  ·  Charge compensation: adding "
           "or removing Li to keep the cell neutral  ·  bars = three doping levels, not repeat runs"),
    drop=[10],
 ),
 10: dict(
    idx=9, out="Slide10_structure_generation.pptx",
    # ⚠ 3×3 을 **표가 아니라 그림**으로 둔다 (1저자 지정 2026-08-18: confusion matrix 양식).
    #   칸 안에 수 + 비율 두 줄, 컬러바, 축 제목 — 표보다 '어느 조합이 두꺼운가' 가 즉시 보인다.
    figs=["docs/figures/seminar/step2_anion_site.png",
          "docs/figures/seminar/step2_site_grid.png"],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Candidate structure generation",
    head="Step 2: Each allowed placement becomes a separate structure",
    sub1="One compound became [b]30 structures[/b] on average, and 3,615 in total.",
    sub2="No site was assumed — every allowed placement was [b]built[/b], and the ranking picked the winner.",
    label="where the anion landed  ·  structures per site pair",
    # 색 정정 (1저자 2026-08-18): 편향 없이 골랐다는 **강점**이므로 파랑이다.
    note="all nine site pairs were populated — [b]the site was never fixed by design[/b]",
    gloss=("Configuration: one specific arrangement of atoms in the cell  ·  "
           "24g / 48h / 4b: Wyckoff labels for the sublattice positions"),
    drop=[10],
 ),
 11: dict(
    idx=10, out="Slide11_lowcost_relaxation.pptx",
    figs=["docs/figures/seminar/step3_survival.png",
          "docs/figures/seminar/step3_screen_axes.png"],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Low-cost structure relaxation",
    head="Step 3: Machine-learned potential screening before DFT",
    sub1="To relax 3,615 structures within budget, a [b]machine-learned potential[/b] was used.",
    sub2="Structures were removed by [b]geometry (volume), not by energy[/b].",
    label="how far each structure moved  ·  the same structures on both screening axes",
    note="[r]100 of 3,615 changed by more than 25 %[/r]; their energies look like the rest.",
    gloss=("Relaxation: moving atoms until the forces vanish  ·  "
           "MLIP: a fast stand-in trained on DFT forces"),
    drop=[10],
 ),
 12: dict(
    idx=11, out="Slide12_representative_structure.pptx",
    figs=["docs/figures/seminar/step4_stability_band.png"],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Representative structure selection",
    # ⛔ 앞 판 문구는 "carried forward" 였다 (1저자: kept 로). 그리고 v6 원본은
    #   용어줄이 헤더 위에 올라타 잘려 있었다 — 이 빌더가 맨 아래로 내린다.
    # ⚠ "per candidate" 는 틀렸다 — 11장에서 candidate = **구조**(3,615개)로 썼다.
    #   "one structure per candidate" 는 "구조당 구조 하나" 로 읽힌다. compound 다.
    head="Step 4: One structure per compound is kept for the full calculations",
    # ⚠ 앞 판 sub1 은 헤더와 같은 말이었다("한 개 남기려고 한 개를 남겼다").
    #   **무엇이 정했는지**로 바꾼다 — 10장 ⛔("제일 안정한 게 이겼다" 는 틀림,
    #   챔피언이 최저 에너지인 경우 9건 중 1건)를 화면에서 미리 막는다.
    sub1="The [b]combined score[/b], not energy alone, picked which configuration was kept.",
    # ⛔⛔ 앞 판 캡션은 "the spread across our three runs ... the runs agree" 였다. **틀렸다.**
    #   세 점은 재현 산포가 아니라 **자리·시드가 다른 세 구조**다 (concentration 열 전부 0.25).
    #   kb/results/site_preference_bar_meaning_2026_08_18.md §4b — 같은 함정 3회째.
    sub2="Each compound produced [b]15 to 150[/b] candidates, so [r]this selection is not a ranking of elements[/r].",
    label="spread of the three top-ranked structures per species",
    # ⚠ em-dash 금지(1/13 덱은 0개). "the choice is real" 도 무슨 선택인지 안 보였다.
    note="median spread [b]0.04 eV / atom[/b]; for the [r]6 of 88 above 0.10[/r], the choice of structure matters",
    # ⚠ 앞 판은 "Convergence" 를 풀어 놨는데 **슬라이드에 그 낱말이 없다**(소불릿에서 빠짐).
    #   화면에 있는 낱말만 푼다.
    gloss=("Representative: the one structure kept for the next steps  ·  "
           "Combined score: several properties folded into one number"),
    drop=[10],
 ),
 13: dict(
    idx=12, out="Slide13_thermal_perturbation.pptx",
    # ⭐ 모식도를 **왼쪽에** 둔다 (1저자 2026-08-19). 앞 판은 히스토그램 하나뿐이라
    #   "왜 막대가 전부 왼쪽인가" 가 결과만 있고 설명이 없었다.
    figs=["docs/figures/seminar/step5_anneal_scheme.png",
          "docs/figures/seminar/step5_anneal_gain.png"],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Thermal perturbation of the selected structure",
    head="Step 5: A short anneal tests whether the arrangement survives",
    sub1="To escape the nearest local minimum, each structure was [b]heated briefly and relaxed again[/b].",
    sub2="The trajectory is [r]not an equilibrium structure and not a conductivity measurement[/r].",
    label="what a short anneal does  ·  the energy it recovered",
    # ⚠ 앞 판 아래 문장("500 K for 50 picoseconds, then relax again — long enough…")은
    #   ① em-dash 가 있었고 ② `a small barri / er` 로 낱말 중간에서 감겼다.
    #   조건은 이제 **왼쪽 모식도가 진다** — 이 줄은 실측만 말한다.
    note="[b]681 of 681[/b] found a lower arrangement, median [b]0.81 eV per cell[/b]; the top ten of the ranking did not change",
    gloss=("Anneal: a brief run at high temperature  ·  "
           "Local minimum: the nearest stable arrangement"),
    drop=[10],
 ),
 14: dict(
    idx=13, out="Slide14_static_pathway.pptx",
    figs=["docs/figures/seminar/step6_li_landscape.png"],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Static lithium transport pathway",
    head="Step 6: The Li energy landscape is mapped on the annealed geometry",
    sub1="To flag transport risk without dynamics, a [b]bond-valence landscape[/b] was computed.",
    sub2="Low-energy valleys are [r]structural pathways, not diffusion coefficients[/r].",
    label="Li landscape on the annealed geometry, undoped versus B₂O₃-doped",
    # ⛔ 출처 명시 (1저자 2026-08-19) — 이 두 판은 **이 발표의 스크리닝 산물이 아니다.**
    #   comp1/modelc DFT 캠페인에서 따로 완주한 예시 두 계다. 앞 판 캡션은 그 말이
    #   없어서 90종 후보 중 둘로 읽혔다.
    note="[r]a worked example from our other DFT study, not a screened candidate[/r]; valleys are low-energy regions, not channels",
    gloss=("Bond-valence map: a cheap estimate of how comfortable an ion is at a point  ·  "
           "Percolation: a low-energy path that spans the crystal"),
    drop=[10],
 ),
 15: dict(
    idx=14, out="Slide15_mechanical_response.pptx",
    # ⭐ 그림 자리를 **비운다** (1저자 2026-08-19) — 1/13 덱의 EOS 두 판
    #   (Li₆PS₅Cl 26.2 GPa · LPSCl1.6 21.7 GPa)을 직접 붙이신다.
    figs=[],
    zone=dict(x=0.55, y=2.44, w=8.90, h=3.28),
    title="Mechanical response of the doped lattice",
    head="Step 7: Stiffness and compressibility are obtained from finite strains",
    # ⛔ 앞 판은 **무엇으로 계산했는지**를 말하지 않았다 (1저자 2026-08-19).
    #   캐스케이드의 02–08 단계는 전부 같은 MLIP 다 (파이프라인 표 07 eos · 08 elastic).
    #   DFT 는 따로 돌린 대조군이고, 90종 값은 **모델 값**이다.
    sub1="Across the pool, [b]the same machine-learned potential[/b] supplied every elastic modulus.",
    sub2="The curves below are [r]DFT on two compositions, not a check of the whole pool[/r].",
    label="equation of state by DFT, undoped Li₆PS₅Cl versus Cl-rich LPSCl1.6",
    note="adding Cl softens the lattice, [b]26.2 to 21.7 GPa[/b]; [r]the pool-wide moduli are potential values, not DFT[/r]",
    gloss=("Equation of state: energy as the cell volume changes  ·  "
           "Stack pressure: the force holding the cell together"),
    drop=[10],
 ),
}


def keep_one_slide(prs, idx):
    """prs 에서 idx 한 장만 남긴다.

    ⚠ python-pptx 에 슬라이드 삭제 API 가 없다. sldIdLst 항목을 지우고 관계도 같이
      끊어야 한다 — 관계를 안 끊으면 파일이 열리긴 하나 고아 파트가 남는다.
    """
    xml_slides = prs.slides._sldIdLst
    keep = list(xml_slides)[idx]
    for sld in list(xml_slides):
        if sld is not keep:
            prs.part.drop_rel(sld.rId)
            xml_slides.remove(sld)
    return prs.slides[0]


def build(n):
    from pptx import Presentation
    from pptx.util import Inches
    import rebuild_cascade_deck as M

    S = SLIDES[n]
    prs = Presentation(SRC)
    sl = keep_one_slide(prs, S["idx"])
    M.clear_zone(sl)                       # 14번 이후(그림·캡션·각주) 전부 정리

    M.set_text(sl.shapes[1], S["title"])
    M.set_text(sl.shapes[5], S["head"])
    M.set_text(sl.shapes[7], S["sub1"])
    M.set_text(sl.shapes[9], S["sub2"])

    # ① 용어줄을 **맨 아래로** — 원래 (5.02, 0.94) 라 헤더 불릿 위에 올라타 잘렸다.
    g = sl.shapes[2]
    g.left, g.top, g.width, g.height = Inches(0.57), Inches(6.86), Inches(8.86), Inches(0.24)
    M.set_text(g, S["gloss"])

    # ② 오른쪽 이탤릭 주석은 지운다 — 그림 밑 이름표와 같은 말이다.
    for i in sorted(S.get("drop", []), reverse=True):
        sl.shapes[i]._element.getparent().remove(sl.shapes[i]._element)

    # figs=[] 이면 그림 자리를 **비운다** — 1저자가 다른 덱의 판을 직접 붙일 때 쓴다
    #   (⚠ add_pics 는 n=0 에서 ZeroDivisionError 다. 여기서 막는다.)
    _figs = S["figs"] if "figs" in S else [S["fig"]]
    if _figs:
        M.add_pics(sl, _figs, zone=S["zone"])
    # 표는 그림 옆/아래에 — 3×3 처럼 **숫자 아홉 개**짜리는 히트맵보다 표가 낫다
    # (1저자 판단 2026-08-18: 파이썬 그림 티가 나느니 덱 표 양식으로).
    if "table" in S:
        t = S["table"]
        M.add_table(sl, t["rows"], t["x"], t["y"], t["w"],
                    size=t.get("size", 11.5), first_col_w=t.get("first_col_w"))
    M.add_fig_label(sl, S["label"], 0.55, 5.96, 8.90)
    M.add_fig_note(sl, S["note"], 0.55, 6.30, 8.90, size=11.5)

    out = os.path.join(ROOT, "docs", S["out"])
    prs.save(out)
    return out


def selftest():
    ok = fail = 0

    def chk(c, m):
        nonlocal ok, fail
        if c:
            ok += 1; print(f"  ✓ {m}")
        else:
            fail += 1; print(f"  ✗ {m}")

    import re as _re
    from pptx import Presentation
    import rebuild_cascade_deck as M

    def plain(t):
        # 강조 표기가 두 벌이다: 덱 [r]/[b] · 대본 [빨강]/[파랑]. 닫는 쪽은 빈 토큰 `[/]`.
        return _re.sub(r"\[/?(?:r|b|빨강|파랑)?\]", "", t)

    chk(os.path.exists(SRC), f"원본 덱이 있다 ({os.path.basename(SRC)})")
    body = plain(open(SCRIPT, encoding="utf-8").read()) if os.path.exists(SCRIPT) else ""
    src = Presentation(SRC)

    for n, S in sorted(SLIDES.items()):
        figs = S["figs"] if "figs" in S else [S["fig"]]
        for f in figs:
            chk(os.path.exists(os.path.join(ROOT, f)),
                f"{n}장 그림이 있다 ({os.path.basename(f)})")
            # 음성 — add_pics 는 repo 루트 기준. 경로가 틀리면 **조용히 건너뛴다**
            chk(f.startswith("docs/"), f"{n}장 그림 경로가 repo 루트 기준")

        # 대본과 1:1 — 헤더·이름표·소불릿까지
        chk(S["head"] in body, f"{n}장 헤더가 대본에도 있다")
        for tag in ("sub1", "sub2"):
            key = plain(S[tag]).split(" — ")[0].split(";")[0].strip()
            chk(key in body, f"{n}장 {tag} 가 대본에도 있다 ('{key[:30]}…')")

        # ★ 글자폭 — 넘치면 두 줄로 감겨 그림을 밀거나 잘린다.
        #   v6 8·9장의 용어줄이 그래서 `of that eleme` 에서 끊겼다.
        for tag, w_in, pt in (("head", 8.49, 17.2), ("sub1", 7.92, 12.8),
                              ("sub2", 7.92, 12.8), ("gloss", 8.86, 8.6),
                              ("note", 8.90, 11.5)):
            w = M._text_pt(plain(S[tag]), pt) / 72.0
            chk(w < w_in - 0.15, f"{n}장 {tag} 한 줄에 들어간다 ({w:.2f} < {w_in} in)")

        # 음성 — 지울 도형이 정말 그 주석인지 (번호로 지우므로 확인이 필수)
        sl = src.slides[S["idx"]]
        for i in S.get("drop", []):
            t = sl.shapes[i].text_frame.text if sl.shapes[i].has_text_frame else ""
            # 지울 것은 **오른쪽 위 이탤릭 주석**이다 — 짧고, 문장이 아니다.
            # 낱말 목록으로 잡으면 장마다 새 낱말이 나온다(‘Method note’). 길이로 본다.
            chk(0 < len(t) < 60 and t.count(".") == 0,
                f"{n}장 지울 shape[{i}] 이 짧은 주석이다 ('{t[:30]}')")
        # 용어줄은 "낱말: 뜻 · 낱말: 뜻" 꼴이다. 낱말을 목록으로 두면 장마다 고쳐야 하므로
        # **모양**으로 본다 — 콜론과 가운뎃점이 있고, 문장이 아니라 정의 나열.
        _t2 = sl.shapes[2].text_frame.text
        chk(":" in _t2 and "·" in _t2 and len(_t2) > 40,
            f"{n}장 옮길 shape[2] 가 용어줄이다 ('{_t2[:34]}…')")

    # 음성 — 영국식 철자가 슬라이드 문장에 섞이면 안 된다 (1저자 지적 2026-08-18)
    allt = " ".join(plain(v) for S in SLIDES.values() for k, v in S.items()
                    if isinstance(v, str))
    brit = [w for w in ("coloured", "colour", "centre", "labelled", "analyse",   # brit-ok: 검출용 단어 목록
                        "behaviour", "normalise") if w in allt.lower()]          # brit-ok: 검출용 단어 목록
    chk(not brit, f"음성: 영국식 철자 없음 ({brit})")

    print(f"\nselftest: {ok} passed, {fail} failed")
    return 1 if fail else 0


if __name__ == "__main__":
    if "--selftest" in sys.argv:
        raise SystemExit(selftest())
    args = [a for a in sys.argv[1:] if a.isdigit()]
    todo = [int(a) for a in args] or sorted(SLIDES)
    for n in todo:
        if n not in SLIDES:
            raise SystemExit(f"⛔ {n}장은 SLIDES 에 없다 (있는 것: {sorted(SLIDES)})")
        print("→", build(n))
