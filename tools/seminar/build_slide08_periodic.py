#!/usr/bin/env python3
"""build_slide08_periodic.py — 8장(주기율표 후보군) 단독 슬라이드를 다시 만든다.

왜 (2026-08-18)
  v6 덱의 8장에 세 가지 문제가 있었다. 전부 **화면에서 눈으로 보이는** 것들이다:
    ① 용어 설명줄(shape 2)이 헤더 불릿(shape 5) 위에 올라타고 오른쪽이 잘렸다
       — `how many compounds of that eleme` 에서 끊겼다.
    ② 같은 단서("superseded snapshot")가 **세 번** 나왔다: 소불릿 · 그림 밑 빨간 줄 ·
       맨 아래 이탤릭. 그 이탤릭도 잘려서 `current rank ing is 0 species` 로 보였다.
    ③ 색이 47종 풀(2026-06-29)이었다. 89종 풀로 바꾸면서 캡션도 같이 가야 한다.

  ⚠ 소불릿의 "superseded snapshot" 은 **변명**이었다. 89종으로 다시 매겨도 상·하위
    집합이 그대로라는 실측이 나왔으므로(카드 아래 참조), 이제 근거 있는 문장으로 바꾼다.

  python3 tools/seminar/build_slide08_periodic.py --selftest
  python3 tools/seminar/build_slide08_periodic.py

이 도구가 못 하는 것
  · 점수를 다시 매기지 않는다. 그림은 plot_seminar_2026_08.fig_periodic 이 만든다.
  · 원본 덱(v6)을 고치지 않는다 — 단독 파일만 만든다. 덱 반영은 사람이 붙여넣는다.
  · 머리말 도형(0–13)의 위치를 바꾸지 않는다. 로고·푸터 규약을 건드리면 안 된다.
"""
import os
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.insert(0, os.path.join(ROOT, "tools", "seminar"))

SRC = os.path.join(ROOT, "docs", "Research_Seminar_2026_08_cascade_story_v6.pptx")
OUT = os.path.join(ROOT, "docs", "Slide08_periodic_table.pptx")
FIGREL = "docs/figures/seminar/roster_periodic_table.png"   # M.add_pics 는 repo 루트 기준
SLIDE_IDX = 7                      # v6 덱의 8장 (0-based)

#: 슬라이드에 인쇄되는 문장 — 대본(kb/seminars/…)의 8장과 **1:1**이어야 한다.
HEAD = "Candidate set (2): 36 cation elements, colored by their best score"
SUB1 = "Late transition metals (Fe–Cu) form the [r]red[/r] block — their oxidation window is the narrowest."
SUB2 = "[b]Group trends[/b] are robust, the exact order is not — the same elements stay at both ends."
LABEL = "best score per element, 89-species pool (2026-08-13)"
NOTE = "[r]no approved ranking yet[/r] — the same elements stayed top and bottom when the set doubled (47 → 89)"
GLOSS = ("Cation: the metal we substitute in  ·  Coverage: how many compounds of that element "
         "were scored  ·  de: stability against the host")


def keep_one_slide(prs, idx):
    """prs 에서 idx 한 장만 남긴다.

    ⚠ python-pptx 에 슬라이드 삭제 API 가 없다. sldIdLst 항목을 지우고 관계도 같이
      끊어야 한다 — 관계를 안 끊으면 파일이 열리긴 하나 고아 파트가 남는다.
    """
    xml_slides = prs.slides._sldIdLst
    keep = list(xml_slides)[idx]
    for sld in list(xml_slides):
        if sld is not keep:
            prs.part.drop_rel(sld.rId)
            xml_slides.remove(sld)
    return prs.slides[0]


def build():
    from pptx import Presentation
    from pptx.util import Inches
    import rebuild_cascade_deck as M

    prs = Presentation(SRC)
    sl = keep_one_slide(prs, SLIDE_IDX)
    M.clear_zone(sl)                       # 14번 이후(그림·캡션·각주) 전부 정리

    M.set_text(sl.shapes[1], "Candidate set on the periodic table")
    M.set_text(sl.shapes[5], HEAD)
    M.set_text(sl.shapes[7], SUB1)
    M.set_text(sl.shapes[9], SUB2)

    # ① 용어줄을 **맨 아래로** 내린다 — 헤더 불릿과 겹치던 자리(5.02, 0.94)를 비운다.
    g = sl.shapes[2]
    g.left, g.top, g.width, g.height = Inches(0.57), Inches(6.86), Inches(8.86), Inches(0.24)
    M.set_text(g, GLOSS)

    # ② 오른쪽 이탤릭 주석(shape 10)은 지운다 — 그림 밑 이름표와 같은 말이다.
    sl.shapes[10]._element.getparent().remove(sl.shapes[10]._element)

    M.add_pics(sl, [FIGREL], zone=dict(x=0.55, y=2.28, w=8.90, h=3.62))
    M.add_fig_label(sl, LABEL, 0.55, 5.96, 8.90)
    M.add_fig_note(sl, NOTE, 0.55, 6.30, 8.90, size=11.5)

    prs.save(OUT)
    return OUT


def selftest():
    ok = fail = 0

    def chk(c, m):
        nonlocal ok, fail
        if c:
            ok += 1; print(f"  ✓ {m}")
        else:
            fail += 1; print(f"  ✗ {m}")

    chk(os.path.exists(SRC), f"원본 덱이 있다 ({os.path.basename(SRC)})")
    fig = os.path.join(ROOT, FIGREL)
    chk(os.path.exists(fig), f"그림이 있다 ({os.path.basename(FIGREL)})")
    # 음성 — add_pics 가 조용히 건너뛰는 경로 실수를 잡는다 (2026-08-18 실제로 그랬다)
    chk(FIGREL.startswith("docs/"), "음성: 그림 경로가 repo 루트 기준이다 (add_pics 규약)")

    # 음성 ① — 대본과 슬라이드 문장이 갈리면 안 된다 (덱과 1:1 규약)
    script = os.path.join(ROOT, "kb", "seminars", "cascade_dopant_screening_story_2026_08.md")
    body = open(script, encoding="utf-8").read() if os.path.exists(script) else ""
    chk(HEAD in body, "헤더 문장이 대본에도 있다")
    chk(LABEL in body, "그림 이름표가 대본에도 있다")
    # ⚠ 소불릿까지 대조한다 — 앞 판은 HEAD/LABEL 만 봐서 SUB2 가 대본과 갈린 채 통과했다.
    # ⚠ 강조 표기가 **두 벌**이다: 덱은 [r]/[b], 대본은 [빨강]/[파랑]. 비교 전에 둘 다 벗긴다
    #   (안 벗기면 같은 문장인데 다르다고 나온다 — 2026-08-18 실제로 그랬다).
    import re as _re

    def _plain(t):
        # 대본의 닫는 표기는 **빈 토큰** `[/]` 다 — 토큰을 선택으로 둬야 잡힌다
        return _re.sub(r"\[/?(?:r|b|빨강|파랑)?\]", "", t)

    body = _plain(body)
    for _tag, _t in (("소불릿1", SUB1), ("소불릿2", SUB2)):
        _key = _plain(_t).split(" — ")[0].strip()
        chk(_key in body, f"{_tag} 이 대본에도 있다 ('{_key[:34]}…')")

    # 음성 ② — 낡은 캡션이 남아 있으면 안 된다 (47종/2026-06-29)
    chk("2026-06-29 snapshot" not in (HEAD + SUB1 + SUB2 + LABEL + NOTE),
        "음성: 47종 시절 캡션이 안 남아 있다")
    chk("superseded" not in SUB2,
        "음성: 소불릿이 '변명'(superseded)이 아니라 근거 문장이다")

    # 음성 ③ — 같은 단서를 두 번 인쇄하지 않는다 (앞 판은 세 번이었다)
    chk(sum(("rank" in t.lower()) for t in (SUB2, NOTE)) <= 2, "순위 단서가 2회 이하")
    # ⚠ 소불릿은 **주장**, 그림 밑 줄은 그 **근거**다. 같은 말을 두 번 쓰면 안 되고,
    #   근거 쪽에 숫자(47 → 89)가 있어야 질문을 막는다.
    chk("both ends" in SUB2 and ("47" in NOTE and "89" in NOTE),
        "소불릿=주장 · 그림밑=근거(47→89 숫자 포함)")

    # ★ 글자폭 — 상자를 넘으면 두 줄로 감겨 아래(그림)를 밀거나 잘린다.
    #   v6 8장의 용어줄이 정확히 그래서 `of that eleme` 에서 끊겼다. 이제 못으로 박는다.
    try:
        import rebuild_cascade_deck as _M
        from pptx import Presentation as _P
        _sl = _P(OUT).slides[0] if os.path.exists(OUT) else None
        if _sl is not None:
            _bad = []
            for _s in _sl.shapes:
                if not _s.has_text_frame or not _s.text_frame.text.strip():
                    continue
                _pt = max((r.font.size.pt for pa in _s.text_frame.paragraphs
                           for r in pa.runs if r.font.size), default=12)
                _w = _M._text_pt(_s.text_frame.text, _pt) / 72.0
                if _w > _s.width / 914400:
                    _bad.append(f"{_s.text_frame.text[:26]}… {_w:.2f}>{_s.width/914400:.2f}in")
            chk(not _bad, f"모든 글상자가 한 줄에 들어간다 ({_bad})")
    except Exception as e:                                   # noqa: BLE001
        chk(False, f"글자폭 검사 실패: {e}")

    # 음성 ④ — 슬라이드에서 지울 도형 번호가 실제로 그 도형인지
    try:
        from pptx import Presentation
        s = Presentation(SRC).slides[SLIDE_IDX]
        chk(len(s.shapes) > 10, "머리말 도형이 10개보다 많다")
        chk("roster" in s.shapes[10].text_frame.text.lower(),
            f"지울 shape[10] 이 정말 그 주석이다 ('{s.shapes[10].text_frame.text[:28]}')")
        chk("Cation" in s.shapes[2].text_frame.text,
            "옮길 shape[2] 가 정말 용어줄이다")
    except Exception as e:                                   # noqa: BLE001
        chk(False, f"원본 덱 도형 검사 실패: {e}")

    print(f"\nselftest: {ok} passed, {fail} failed")
    return 1 if fail else 0


if __name__ == "__main__":
    if "--selftest" in sys.argv:
        raise SystemExit(selftest())
    print("→", build())
