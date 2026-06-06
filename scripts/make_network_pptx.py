#!/usr/bin/env python3
"""Figure 1(b) — network-solver schematic as an EDITABLE .pptx.

Every particle is a real PowerPoint OVAL; every resistor is a real
zigzag FREEFORM line.  Edge colors are COMPUTED from the two endpoint
particle phases (yellow SE-SE / blue AM-SE / red AM-AM) so they cannot
be wrong.  Open in PowerPoint and drag / recolor any shape by hand.

Run:  python3 scripts/make_network_pptx.py
Out:  docs/figures/network_solver_schematic.pptx
"""
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

rng = np.random.default_rng(3)

# ── colors ──
def C(h): return RGBColor.from_string(h)
C_SE, C_SE_E = C('F6C623'), C('C08A0A')
C_AM, C_AM_E = C('9B9B9B'), C('3F3F3F')
E_SESE, E_AMSE, E_AMAM = C('F3BF1E'), C('3F7FD0'), C('E22B2B')
BB_Y, BB_R = C('F5A800'), C('EE1111')
SUS_FC, SUS_EC = C('E8807F'), C('C64A4A')
BULK_FC, BULK_EC = C('E8D3A0'), C('B89030')

# ── data-space geometry (same as the matplotlib version) ──
X0, X1 = 0.5, 9.5
SUS_Y, BULK_Y = 10.0, 1.0
SE_TOP, AM_BOT = 8.3, 2.7

nodes = []
n_se = 9
se_y = np.linspace(BULK_Y+0.15, SE_TOP, n_se)
se_x = 3.0 + 0.50*np.sin(np.linspace(0.2, 3.4, n_se))
se_chain = list(zip(se_x, se_y))
for (x, y) in se_chain: nodes.append((x, y, 'SE', 0.20, True))
SE_GAP_X = se_chain[-1][0]

n_am = 9
am_y = np.linspace(SUS_Y-0.15, AM_BOT, n_am)
am_x = 6.7 + 0.40*np.sin(np.linspace(0.0, 3.0, n_am))
am_chain = list(zip(am_x, am_y))
for (x, y) in am_chain: nodes.append((x, y, 'AM', 0.30, True))
AM_GAP_X = am_chain[-1][0]

def clash(x, y, r, pad=0.05):
    for (nx, ny, ph, nr, bb) in nodes:
        if (x-nx)**2 + (y-ny)**2 < (r+nr+pad)**2: return True
    return False

gx = np.linspace(X0+0.4, X1-0.4, 12)
gy = np.linspace(BULK_Y+0.35, SUS_Y-0.35, 13)
for y in gy:
    for x in gx:
        xx = x + rng.uniform(-0.16, 0.16); yy = y + rng.uniform(-0.16, 0.16)
        if yy > SE_TOP and abs(xx - SE_GAP_X) < 0.75: continue
        if yy < AM_BOT and abs(xx - AM_GAP_X) < 0.75: continue
        if yy > SE_TOP:   is_am = True
        elif yy < AM_BOT: is_am = False
        else:             is_am = rng.random() < 0.13
        r = rng.uniform(0.26, 0.34) if is_am else 0.16
        if clash(xx, yy, r): continue
        nodes.append((xx, yy, 'AM' if is_am else 'SE', r, False))

def spring(p, q, amp=0.062, period=0.30):
    p = np.array(p, float); q = np.array(q, float)
    v = q - p; L = np.hypot(*v)
    if L < 1e-6: return [tuple(p)]
    u = v / L; perp = np.array([-u[1], u[0]])
    lead = 0.18
    n_teeth = max(3, int(round((L*(1-2*lead))/period)))
    ts = np.linspace(lead, 1-lead, n_teeth*2+1)
    pts = [tuple(p)]
    for k, t in enumerate(ts):
        off = 0.0 if (k == 0 or k == len(ts)-1) else amp*(1 if k % 2 else -1)
        pts.append(tuple(p + v*t + perp*off))
    pts.append(tuple(q))
    return pts

SE = {'SE'}
def near(a, b, d=1.15): return (a[0]-b[0])**2 + (a[1]-b[1])**2 < d*d
edges = {'AM-AM': [], 'AM-SE': [], 'SE-SE': []}
for i in range(len(nodes)):
    for k in range(i+1, len(nodes)):
        a, b = nodes[i], nodes[k]
        if not near(a, b): continue
        s1 = a[2] in SE; s2 = b[2] in SE
        if s1 and s2:               edges['SE-SE'].append((a, b))
        elif (not s1) and (not s2): edges['AM-AM'].append((a, b))
        else:                       edges['AM-SE'].append((a, b))

# ── data-coord → slide EMU (equal x/y scale so circles stay round) ──
S = 0.62          # inches per data unit
MX, MY = 0.35, 0.30
def sx(x): return MX + x*S
def sy(y): return MY + (11 - y)*S          # flip: data-y up, slide-y down
EMU = 914400
def E(inch): return Emu(int(round(inch*EMU)))

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
shapes = slide.shapes

def add_bar(y_top_data, y_bot_data, fc, ec):
    left = E(sx(X0-0.2)); top = E(sy(y_top_data))
    w = E((X1-X0+0.4)*S); h = E((y_top_data-y_bot_data)*S)
    sh = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = ec; sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    return sh

add_bar(SUS_Y+0.9, SUS_Y, SUS_FC, SUS_EC)        # SUS collector (top)
add_bar(BULK_Y, BULK_Y-0.9, BULK_FC, BULK_EC)    # bulk reservoir (bottom)

def add_spring(pts, color, width_pt):
    flat = [(E(sx(px)), E(sy(py))) for (px, py) in pts]
    fb = shapes.build_freeform(flat[0][0], flat[0][1], scale=1.0)
    fb.add_line_segments(flat[1:], close=False)
    sh = fb.convert_to_shape()
    sh.fill.background()
    sh.line.color.rgb = color; sh.line.width = Pt(width_pt)
    sh.shadow.inherit = False
    return sh

# edges (red under blue under yellow, same as raster)
for a, b in edges['AM-AM']: add_spring(spring((a[0],a[1]),(b[0],b[1])), E_AMAM, 1.5)
for a, b in edges['AM-SE']: add_spring(spring((a[0],a[1]),(b[0],b[1])), E_AMSE, 1.4)
for a, b in edges['SE-SE']: add_spring(spring((a[0],a[1]),(b[0],b[1])), E_SESE, 1.3)

# backbones (thicker bold springs over the chain)
for chain, col in [(se_chain, BB_Y), (am_chain, BB_R)]:
    for m in range(len(chain)-1):
        add_spring(spring(chain[m], chain[m+1]), col, 3.2)

# nodes (ovals on top)
def add_node(x, y, ph, r, bb):
    fc, ec = (C_SE, C_SE_E) if ph == 'SE' else (C_AM, C_AM_E)
    d = E(2*r*S)
    left = E(sx(x) - r*S); top = E(sy(y) - r*S)
    sh = shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = ec; sh.line.width = Pt(1.6 if bb else 1.0)
    sh.shadow.inherit = False
    return sh

for (x, y, ph, r, bb) in nodes: add_node(x, y, ph, r, bb)

# X break marks (two crossing line segments each, in the empty pockets)
def add_X(xc, yc, color, half=0.28):
    for dx, dy in [((-half, -half), (half, half)), ((-half, half), (half, -half))]:
        add_spring([(xc+dx[0], yc+dx[1]), (xc+dy[0], yc+dy[1])], color, 4.0)
sx_, sy_ = se_chain[-1]; add_X(sx_, sy_+0.75, BB_Y)
ax_, ay_ = am_chain[-1]; add_X(ax_, ay_-0.75, BB_R)

# ── legend (right) ──
LX, LY, LW, LH = 9.0, 0.5, 4.0, 6.6   # inches
box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(LX), Inches(LY), Inches(LW), Inches(LH))
box.fill.solid(); box.fill.fore_color.rgb = C('FFFFFF')
box.line.color.rgb = C('333333'); box.line.width = Pt(1.4)
box.shadow.inherit = False

def txt(x, y, s, size=12, bold=False):
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(3.6), Inches(0.35))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = s
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = C('222222')

def leg_oval(x, y, fc, ec, d=0.26):
    sh = shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = ec; sh.line.width = Pt(1.2); sh.shadow.inherit = False

def leg_zig(x, y, color, w=2.0):
    pts = [(x, y), (x+0.12, y-0.07), (x+0.24, y+0.07),
           (x+0.36, y-0.07), (x+0.48, y+0.07), (x+0.60, y)]
    flat = [(E(px), E(py)) for (px, py) in pts]
    fb = shapes.build_freeform(flat[0][0], flat[0][1], scale=1.0)
    fb.add_line_segments(flat[1:], close=False)
    sh = fb.convert_to_shape(); sh.fill.background()
    sh.line.color.rgb = color; sh.line.width = Pt(w); sh.shadow.inherit = False

txt(LX+0.25, LY+0.15, 'NODES', 13, True)
leg_oval(LX+0.35, LY+0.65, C_SE, C_SE_E);       txt(LX+0.9, LY+0.62, 'yellow = SE (ionic, majority matrix)', 11)
leg_oval(LX+0.32, LY+1.12, C_AM, C_AM_E, 0.33); txt(LX+0.9, LY+1.12, 'gray = AM (electronic, minority islands)', 11)
txt(LX+0.25, LY+1.75, 'CONTACTS (resistors)', 13, True)
leg_zig(LX+0.35, LY+2.30, E_SESE); txt(LX+1.15, LY+2.18, 'yellow = SE-SE', 11)
leg_zig(LX+0.35, LY+2.72, E_AMSE); txt(LX+1.15, LY+2.60, 'blue  = AM-SE', 11)
leg_zig(LX+0.35, LY+3.14, E_AMAM); txt(LX+1.15, LY+3.02, 'red   = AM-AM', 11)
txt(LX+0.25, LY+3.75, 'BACKBONES', 13, True)
leg_zig(LX+0.35, LY+4.30, BB_Y, 3.0); txt(LX+1.15, LY+4.18, 'yellow → bulk (Li+ path); X = no SE to SUS', 10)
leg_zig(LX+0.35, LY+4.78, BB_R, 3.0); txt(LX+1.15, LY+4.66, 'red → SUS (e- path); X = no AM to bulk', 10)

import os; os.makedirs('docs/figures', exist_ok=True)
out = 'docs/figures/network_solver_schematic.pptx'
prs.save(out)
print(f"saved: {out}  (AM-AM={len(edges['AM-AM'])}, AM-SE={len(edges['AM-SE'])}, "
      f"SE-SE={len(edges['SE-SE'])}, nodes={len(nodes)})")
